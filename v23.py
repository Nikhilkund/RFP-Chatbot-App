import streamlit as st
from pathlib import Path
import os
import tempfile
from io import BytesIO
import docx
from PyPDF2 import PdfReader
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from uuid import uuid4
from dotenv import load_dotenv
import json
from docx.shared import Inches as DocxInches, Pt as DocxPt, RGBColor as DocxRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.style import WD_STYLE_TYPE
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_ALIGN_VERTICAL
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from datetime import date
import re
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

__import__('pysqlite3')
import sys
sys.modules['sqlite3'] = sys.modules.pop('pysqlite3')
# --- End sqlite3 fix ---

# The rest of your existing imports would follow immediately after this block.
# For example, your next line would be:
from pathlib import Path
# --- Import the Groq client directly ---
from groq import Groq


# Langchain imports
from langchain.chains import RetrievalQAWithSourcesChain
from langchain.chains.qa_with_sources.loading import load_qa_with_sources_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma
from langchain_groq import ChatGroq
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document
from langchain.prompts import PromptTemplate

# New loaders for robust document parsing
from langchain_community.document_loaders import UnstructuredPowerPointLoader
from langchain_community.document_loaders import UnstructuredWordDocumentLoader


# --- Configuration and Global Variables ---
load_dotenv()

# Constants
CHUNK_SIZE = 1000
SLIDE_CHUNK_SIZE = 700
SLIDE_CHUNK_OVERLAP = 100

EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
VECTORSTORE_DIR = Path(__file__).parent / "resources/vectorstore"
COLLECTION_NAME = "real_estate"
LOGO_PATH = "logo.png"


llm = None
vector_store = None

# Removed "Overall Summary", "RFP Overview Details", "Cost Proposal Form", "Scorecard Evaluation"
NEW_SECTION_TITLES = [
    "Prospect & RFP Background",
    "Scope of Work Details",
    "Service Level Agreements (SLAs)",
    "RFP Submission Information",
    "RFP Schedule",
    "RFP Evaluation Criteria",
    "SWOT Analysis" # Only SWOT remains from the previous tail
]

PROMPT = PromptTemplate.from_template(
    """Given the following extracted parts of a long document and a question, create a final answer with sources.
If you don't know the answer, just say that you don't know. Don't try to make up an answer.
ALWAYS return a "SOURCES" part in your answer.

QUESTION: {question}
=========
{summaries}
=========
FINAL ANSWER:"""
)

EXAMPLE_PROMPT = PromptTemplate.from_template(
    """CONTENT: {page_content}
SOURCE: {source}"""
)

def initialize_components(model_name_for_init="deepseek-r1-distill-llama-70b"):
    global llm, vector_store

    if llm is None or llm.model_name != model_name_for_init:
        llm = ChatGroq(model=model_name_for_init, temperature=0.5, max_tokens=500)
    
    if vector_store is None:
        ef = HuggingFaceEmbeddings(
            model_name=EMBEDDING_MODEL,
            model_kwargs={"trust_remote_code": True}
        )
        vector_store = Chroma(
            collection_name=COLLECTION_NAME,
            embedding_function=ef,
            persist_directory=str(VECTORSTORE_DIR)
        )


def extract_text_from_pdf(file_path):
    text = ""
    try:
        with open(file_path, "rb") as file:
            reader = PdfReader(file)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
    except Exception as e:
        st.error(f"Error reading PDF: {e}")
        return ""
    return text


def extract_text_from_docx_unstructured(file_path):
    """Extracts text from a DOCX file using UnstructuredWordDocumentLoader."""
    try:
        loader = UnstructuredWordDocumentLoader(file_path)
        docs = loader.load()
        full_text = "\n".join([doc.page_content for doc in docs])
        return full_text
    except Exception as e:
        st.error(f"Error extracting text from DOCX using Unstructured: {e}")
        return ""


def extract_text_from_pptx_unstructured(file_path):
    """Extracts text from a PPTX file using UnstructuredPowerPointLoader with OCR."""
    try:
        loader = UnstructuredPowerPointLoader(file_path, ocr_strategy="hi_res")
        docs = loader.load()
        full_text = "\n".join([doc.page_content for doc in docs])
        return full_text
    except Exception as e:
        st.error(f"Error extracting text from PPTX using Unstructured with OCR: {e}")
        st.info("Ensure Tesseract OCR engine is installed and in your system's PATH.")
        return ""


def process_document_and_update_state(file_path, file_type, original_file_name):
    """
    This function orchestrates text extraction and vector store update.
    It directly updates session state variables.
    """
    global vector_store

    st.session_state.processed_doc_content = None
    st.session_state.current_uploaded_file_name = None
    st.session_state.processing_done = False
    st.session_state.doc_analysis_data = {}  # type: ignore
    st.session_state.understanding_doc_analysis_data = {} # Initialize for new analysis


    progress_bar = st.progress(0)
    status_text_placeholder = st.empty()

    expected_steps = 6
    current_step = 0

    try:
        status_text_placeholder.text("Initializing Components...")
        initialize_components("deepseek-r1-distill-llama-70b")  
        current_step += 1
        progress_bar.progress(current_step / expected_steps)

        status_text_placeholder.text("Resetting vector store...✅")
        if vector_store:
            vector_store.delete_collection()
            ef = HuggingFaceEmbeddings(
                model_name=EMBEDDING_MODEL,
                model_kwargs={"trust_remote_code": True}
            )
            vector_store = Chroma(
                collection_name=COLLECTION_NAME,
                embedding_function=ef,
                persist_directory=str(VECTORSTORE_DIR)
            )
        else:
            st.warning("Vector store not initialized, skipping reset. This should not happen if initialize_components runs first.")
        current_step += 1
        progress_bar.progress(current_step / expected_steps)

        status_text_placeholder.text(f"Extracting text from {file_type} document...✅")
        extracted_text = ""
        if file_type == "pdf":
            extracted_text = extract_text_from_pdf(file_path)
        elif file_type == "docx":
            extracted_text = extract_text_from_docx_unstructured(file_path)
        elif file_type == "pptx":
            extracted_text = extract_text_from_pptx_unstructured(file_path)
        else:
            st.error(f"Unsupported file type: {file_type} ❌. Please upload a PDF, DOCX, or PPTX.")
            return False

        if not extracted_text.strip():
            st.warning("No readable text found in the document. Please check the file content. ❌")
            return False
        
        current_step += 1
        progress_bar.progress(current_step / expected_steps)

        status_text_placeholder.text("Splitting text into chunks for RAG...✅")
        text_splitter = RecursiveCharacterTextSplitter(
            separators=["\n\n", "\n", ".", " "],
            chunk_size=CHUNK_SIZE
        )
        if isinstance(extracted_text, str):
            doc_to_split = [Document(page_content=extracted_text, metadata={"source": original_file_name})]
        else:  
            doc_to_split = extracted_text  
            for doc in doc_to_split:
                if "source" not in doc.metadata:
                    doc.metadata["source"] = original_file_name

        docs = text_splitter.split_documents(doc_to_split)

        if not docs:
            st.warning("Document content too short to generate meaningful chunks for the vector database. ❌")
            return False
        
        current_step += 1
        progress_bar.progress(current_step / expected_steps)

        status_text_placeholder.text(f"Adding {len(docs)} chunks to vector database...✅")
        uuids = [str(uuid4()) for _ in range(len(docs))]
        vector_store.add_documents(docs, ids=uuids)
        
        current_step += 1
        progress_bar.progress(current_step / expected_steps)

        status_text_placeholder.text("Done adding docs to vector database! You can now chat. ✅")
        st.session_state.processed_doc_content = extracted_text
        st.session_state.current_uploaded_file_name = original_file_name
        st.session_state.processing_done = True
        current_step += 1
        progress_bar.progress(current_step / expected_steps)

        return True

    except Exception as e:
        st.error(f"An error occurred during document processing: {e} ❌")
        return False
    finally:
        progress_bar.empty()
        status_text_placeholder.empty()


def generate_answer(query):
    """
    Generates an answer based on the user's query and the current vector store.
    """
    if vector_store is None:
        raise RuntimeError("Vector database is not initialized. Please upload a document first.")
    
    initialize_components("llama-3.3-70b-versatile")

    qa_chain = load_qa_with_sources_chain(llm, chain_type="stuff",
                                         prompt=PROMPT,
                                         document_prompt=EXAMPLE_PROMPT)
    chain = RetrievalQAWithSourcesChain(combine_documents_chain=qa_chain, retriever=vector_store.as_retriever(),
                                        reduce_k_below_max_tokens=True, max_tokens_limit=8000,
                                        return_source_documents=True)
    
    result = chain.invoke({"question": query}, return_only_outputs=True)
    
    sources_docs = [doc.metadata.get('source', 'Unknown') for doc in result.get('source_documents', [])]
    
    return result['answer'], sources_docs

# --- Functions for Document Generation (PPTX/DOCX based on original prompt) ---

def analyze_with_llm_for_docs(context_text, model_name_for_analysis="deepseek-r1-distill-llama-70b"):
    """
    Uses the LLM to analyze the document and extract structured content
    for PPTX and DOCX generation.
    """
    client = Groq(api_key=os.getenv("GROQ_API_KEY"))

    # --- Company Information to be injected into the prompt ---
    company_info = """
BCT Company Overview
BCT provides technology-driven solutions for organizations and governments, focusing on efficiency, sustainability, and security.
Key Service Areas:
Managed IT & Infrastructure: 24/7 helpdesk, infrastructure & security services for asset protection, managed application and backend support, and AI-driven retail store operations.
Smart Governance & Public Sector: Solutions to improve citizen-government relationships, including citizen engagement management, billing & payment systems, AI-driven governance, business automation, e-Procurement, and smart/preventive asset maintenance.
Smart Healthcare: AI-driven public health solutions like remote patient monitoring, population health analytics, track & trace, and smart asset tracking for critical medical resources.
Smart Cities & Infrastructure: Solutions to enhance urban quality of life, including Visual AI for surveillance, smart energy management (meter monitoring, grid stability), connected parks, smart public restrooms, smart transportation/vehicle telematics, and smart classrooms/e-learning platforms.
Sustainability: Solutions are integrated across all projects to reduce environmental footprints and improve operational performance.
Global Capability Center (GCC) Services
BCT supports the setup and optimization of GCCs for clients in BFSI, Retail, Healthcare & Lifescience, and the public sector.
ISG Provider Lens™ Report (2025) Positioning: "Contender" in both "GCC Design and Setup" and "Optimization and Enhancement" quadrants.
Target Clients ("Sweet Spot"): Organizations seeking a strategic partner for co-creation, scaling, and optimization of both new (greenfield) and existing GCCs, with a focus on operational rigor and digital enablement.
Key GCC Capabilities & Approach:
End-to-End Lifecycle Management: Provides comprehensive support including strategic planning, location analysis, regulatory readiness, setup, scaling, and ongoing operations management using agile delivery.
Digital-First & Human-Centric Design: Integrates AI, ML, automation, and low-code platforms into workflows while focusing on employee experience (EX), engagement, and continuous learning.
Co-Managed & Flexible Models: Offers tailored, co-managed service models with robust SLAs, scalable talent ramp-up, and governance-driven delivery that provides clients with transparency, control, and innovation.
Strong Partner Ecosystem: Leverages a curated network of technology, cloud, cybersecurity, and infrastructure partners.
Proven Project Outcomes:
Reduced IT operations TCO by 80% for a large global bookseller.
Improved workforce efficiency by 20% for a Southeast Asian government by streamlining 140 services across 35 departments.
Reduced infrastructure costs by 20% for a leading tax, assurance, and advisory firm.
Enhanced capabilities on the Temenos platform for BFSI clients via the acquisition of Stoics IT.
Technologies & Skills
Core Technologies Utilized:
Artificial Intelligence (AI) & Machine Learning (ML)
Automation & Low-code platforms
Predictive Analytics & Data Analytics
Cloud Technologies & Cloud-Native Architectures
Cybersecurity
Visual AI, Vehicle Telematics
Temenos Platform
Future Roadmap & Emerging Tech Focus:

Technology Expansion: Scaling offerings in AI governance, agentic AI, AR/VR, and advanced computing.
Geographic Expansion: Planned expansion to Mexico City, London, Dubai, and various cities in India.
M&A Strategy: Continuing acquisitions to expand technology and service capabilities.
Industry Deepening: Developing more industry-specific digital solutions for core sectors.
Core Company Skills:

Strategic & Advisory: Consulting, location analysis, and operating model design.
Technical & Digital: Expertise in AI, automation, cloud, and data analytics.
Domain Knowledge: Industry-specific expertise for tailored solutions (CoEs).
Talent Management: Sourcing, onboarding, and upskilling specialized talent.
Change Management & Governance: Managing compliance and organizational transitions.
"""

    # --- UPDATED PROMPT TO INCLUDE "RFP Number" ---
    prompt = f"""Analyze this RFP document and extract content organized into the following detailed sections.
    For each section, provide comprehensive and detailed information.

    - **Prospect & RFP Background**: Detailed information about the client, the specific problem the RFP aims to solve, the strategic context, and any historical information. Extract as much relevant context as possible. Output as a JSON dictionary with specific keys like 'Client Name', 'Project Name', 'RFP Number', 'Industry', 'Project Purpose', 'Background', 'Current Challenges', 'Service Agreement Details', 'Key Contact Name', 'Contact Email', 'Contact Phone', 'Virtual Meeting Details'.
      For 'Service Agreement Details', **meticulously analyze the agreement, contract, or terms sections in the document. Determine the duration of the agreement in years (e.g., "5 years"). If it's in months (e.g., "36 months"), convert it to years (e.g., "3 years"). Also, search for any 'LBE Subcontract' requirements and provide a concise, one-line summary (e.g., "LBE subcontract required for 20% of work"). Combine these into a single string for 'Service Agreement Details' (e.g., "5 years agreement; LBE subcontract required for 20% of work"). If only one detail is found, provide just that detail.**
      For 'Virtual Meeting Details', provide a nested dictionary with 'Meeting Topic', 'Meeting ID', 'Passcode', 'Join Link'. **Only include a sub-key within 'Virtual Meeting Details' if its value is explicitly found and is not empty or a generic placeholder.**
      **Only include a top-level key in the JSON output if its value is explicitly found in the document and is not empty, 'N/A', 'Not Specified', or a generic example/placeholder.** If a key is not found, omit it from the JSON output.
      Example JSON value: {{"Client Name": "Alabama Industrial Development Training (AIDT)", "Project Name": "Utility Management System Implementation", "RFP Number": "S24094", "Industry": "Workforce Development", "Service Agreement Details": "5 years agreement; LBE subcontract required for 20% of work", "Key Contact Name": "John Doe", "Contact Email": "john.doe@example.com", "Contact Phone": "555-123-4567", "Virtual Meeting Details": {{"Meeting Topic": "Pre-bid Conference", "Meeting ID": "123-456-7890", "Passcode": "RFP2025", "Join Link": "https://actual-meeting-link.com/join"}}}}
    - **Scope of Work Details**: A very detailed breakdown of all required tasks, deliverables, specific services, responsibilities of both parties, and any technologies or methodologies mentioned. For each task or area, provide a detailed explanation of what is required. Structure this as a JSON dictionary where each key is a major work area (e.g., "Technical Upgrade", "Migration Steps", "Testing") and its value is a very detailed string description or a JSON list of concise bullet points for that area. Do NOT nest dictionaries for sub-areas.
      Example JSON value: {{"Technical Upgrade": ["Detailed explanation of database upgrade process, including versions and migration steps.", "In-depth description of customization retrofitting process."], "Testing": ["Comprehensive plan for user acceptance testing (UAT) and system integration testing (SIT), detailing responsibilities."]}}
    - **Service Level Agreements (SLAs)**: Analyze the document for any explicit Service Level Agreements. For each SLA found, describe the *existing SLA/client's current requirement* and then outline the *vendor's responsibilities or what needs to be done* to meet this SLA. Output as a JSON list of dictionaries, where each dictionary has keys like "Metric", "Client's Current", and "Vendor's Responsibility". **Only include a dictionary in the list if the SLA is explicitly defined.**
      Example JSON value: [{{ "Metric": "System Uptime", "Client's Current": "99.5% availability", "Vendor's Responsibility": "Maintain 99.9% uptime, implement proactive monitoring and immediate response." }}]
    - **RFP Submission Information**: Provide comprehensive details about submitting the proposal packet. **Crucially, identify who the proposal should be submitted to, including their name, title, or department.** Output as a JSON dictionary with keys: 'Proposal Packet Submission Info', 'Required Forms', 'Submission Method', 'Submission Contact/Whom to Submit'. Only include a key if its value is explicitly found and not empty/placeholder.
      Example JSON value: {{"Proposal Packet Submission Info": "Submit via online portal by 5 PM EST", "Required Forms": ["Form A", "Form B"], "Submission Method": "Online Portal", "Submission Contact/Whom to Submit": "Procurement Department, Attn: Jane Smith"}}
    - **RFP Schedule**: Extract all critical dates, deadlines, **timelines, and high-level milestones**. This includes issue date, pre-bid conference, Q&A deadlines, proposal submission deadline, evaluation period, award date, and contract start date. Output as a JSON dictionary. Only include a key if its date is explicitly found.
      Example JSON value: {{"Issue Date": "January 30, 2025", "Proposal Deadline": "March 14, 2025", "High-Level Milestones": "Project Kick-off: May 1, 2025; Phase 1 Go-live: August 1, 2025"}}
    - **RFP Evaluation Criteria**: Provide details on evaluation criteria, including Minimum Qualification, Evaluation Criteria, and any Labour/Staffing Description details. Output as a JSON dictionary with keys: 'Minimum Qualifications', 'Evaluation Criteria', 'Labor/Staffing Description'. For 'Evaluation Criteria', provide a list of key criteria. For 'Labor/Staffing Description', summarize any specific requirements for personnel (e.g., certifications, experience levels, team composition). Only include a key if its value is explicitly found.
      Example JSON value: {{"Minimum Qualifications": "5 years experience in Oracle EBS", "Evaluation Criteria": ["Technical Approach", "Past Performance", "Cost"], "Labor/Staffing Description": "Project Manager must have PMP certification; 2 Oracle DBAs required."}}
    - **SWOT Analysis**: Perform a SWOT analysis from the perspective of a vendor bidding on this RFP.
      - **Strengths**: Based on the provided company information, identify internal strengths that are **directly relevant to meeting the specific requirements of this RFP**.
      - **Weaknesses**: Based on the provided company information, identify internal weaknesses or gaps that could be a **disadvantage for this specific RFP project**.
      - **Opportunities**: Based on the **RFP document content**, identify external factors or project aspects that the vendor could leverage for additional value or success.
      - **Threats**: Based on the **RFP document content**, identify external risks, competitive factors, or challenges mentioned that could impact the project.
      Company Information:
      {company_info}
      Provide a JSON dictionary with keys 'Strengths', 'Weaknesses', 'Opportunities', 'Threats', each containing a JSON list of concise bullet points (strings).
      Example JSON value: {{"Strengths": ["Proven experience in AI-driven governance matches project needs."], "Weaknesses": ["Limited presence in the client's geographic region."], "Opportunities": ["The RFP mentions a future phase 2, offering long-term engagement."], "Threats": ["The RFP requires a quick turnaround, which could strain resources."]}}

    Respond ONLY with JSON where each key is a section title from the list above, and the value is its corresponding content extracted from the document.
    Keep technical details intact. Omit sections with no relevant content (do not include the key if no content).

    Document Content:
    {context_text[:8000]}"""

    try:
        response = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model=model_name_for_analysis,
            response_format={"type": "json_object"},
            temperature=0.3
        )
        return json.loads(response.choices[0].message.content)
    except json.JSONDecodeError as e:
        st.error(f"Error decoding JSON response from LLM: {e}")
        st.error(f"Raw response: {response.choices[0].message.content}")
        return {}
    except Exception as e:
        st.error(f"An unexpected error occurred during LLM document analysis: {str(e)}")
        if "GROQ_API_KEY" not in os.environ:
            st.error("GROQ_API_KEY environment variable not found. Please set it in your .env file or Streamlit secrets.")
        return {}

def generate_new_pptx_option1(analysis_data, title_bg_info, second_page_bg_info, content_bg_info):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[-1]
    
    # Function to apply background image and set default color
    def apply_background_image_pptx(slide, image_bytes, image_name, prs_obj, default_color_rgb):
        prs_width = prs_obj.slide_width
        prs_height = prs_obj.slide_height

        if image_bytes:
            try:
                with tempfile.TemporaryDirectory() as tmpdir:
                    temp_image_path = Path(tmpdir) / image_name
                    with open(temp_image_path, "wb") as f:
                        f.write(image_bytes)
                    
                    slide.background.fill.background()
                    slide.shapes.add_picture(str(temp_image_path), 0, 0, prs_width, prs_height)  
            except Exception as e:
                st.warning(f"Could not set background image ({image_name}): {e}. Slide will have default solid background.")
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = default_color_rgb
        else:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = default_color_rgb


    # --- First Slide (Title Page) ---
    title_slide = prs.slides.add_slide(blank_slide_layout)
    # Apply background image, default to white for the first slide if no custom image
    apply_background_image_pptx(title_slide, title_bg_info["bytes"], title_bg_info["name"], prs, RGBColor(255, 255, 255))

    # Extract dynamic information for the title slide
    client_name_val = "N/A" # Default
    rfp_number_val = "N/A" # Default
    # --- CHANGE 1: MODIFIED DEFAULT RFP NAME ---
    rfp_name_val = "Project Name" # Default to "Project Name" instead of "N/A"
    
    # Try to get Client Name - SAFELY ACCESS THE DICTIONARY
    bg_data = analysis_data.get("Prospect & RFP Background", {}) # Use .get() to avoid KeyError
    if isinstance(bg_data, dict) and "Client Name" in bg_data:
        client_name_val = bg_data["Client Name"]
    elif isinstance(bg_data, str): # Attempt to parse from string if LLM returned string
        match = re.search(r"(?:Client Name|Client):\s*(.*?)(?:\n|$)", bg_data, re.IGNORECASE)
        if match:
            client_name_val = match.group(1).strip()
        else: # Fallback to trying to get first non-empty line
            lines = [line.strip() for line in bg_data.split('\n') if line.strip()]
            if lines:
                client_name_val = lines[0]


    # No longer pulling from "RFP Overview Details" as it's removed
    # If RFP Number/Name are critical for the title slide, you'd need to adjust the LLM prompt for "Prospect & RFP Background"
    # to include them, or manually extract from the general text if possible.
    # For now, they remain "N/A" unless explicitly found in "Prospect & RFP Background" as a string
    if isinstance(bg_data, str):
        match_num = re.search(r"(?:RFP Number|RFQ Number):\s*(.*?)(?:\n|$)", bg_data, re.IGNORECASE)
        if match_num:
            rfp_number_val = match_num.group(1).strip()
        
        match_name = re.search(r"(?:RFP Name|RFQ Name|Project Name|Project Title):\s*(.*?)(?:\n|$)", bg_data, re.IGNORECASE)
        if match_name:
            # Only update if a valid name is found, otherwise keep the "Project Name" default
            found_name = match_name.group(1).strip()
            if found_name:
                rfp_name_val = found_name


    current_date_val = date.today().strftime("%B %d, %Y") # Format: June 05, 2025

    # Position the text elements for the title slide
    # Create a single text box and add paragraphs to it for precise spacing and alignment
    
    left_align_pos = Inches(1.0) # Left margin for content
    top_start = Inches(2.5) # Starting height for the first line
    width_box = prs.slide_width - (2 * left_align_pos) # Box spans most of the width
    height_box = Inches(5.0) # Ample height for all lines and spaces

    content_box = title_slide.shapes.add_textbox(left_align_pos, top_start, width_box, height_box)
    text_frame = content_box.text_frame
    text_frame.clear() # Clear any default text

    # Client Name (22pt, white, no key)
    p_client = text_frame.add_paragraph()
    p_client.text = client_name_val
    p_client.font.size = Pt(22)
    p_client.font.bold = True
    p_client.font.color.rgb = RGBColor(255, 255, 255) # White text
    p_client.alignment = PP_ALIGN.LEFT

    # RFP Number: [number] (18pt, white, key-value pair)
    p_number_kv = text_frame.add_paragraph()
    p_number_kv.text = f"RFP Number: {rfp_number_val}" # Explicit key-value pair
    p_number_kv.font.size = Pt(18)
    p_number_kv.font.color.rgb = RGBColor(255, 255, 255) # White text
    p_number_kv.alignment = PP_ALIGN.LEFT

    # Two blank lines (empty paragraphs) - Use enough runs to visually represent spacing
    for _ in range(2):
        blank_p = text_frame.add_paragraph()
        blank_p.text = " " # Add a space to ensure the paragraph takes up height
        blank_p.font.size = Pt(18) # Set font size for consistent spacing


    # RFP Name (18pt, white, no key) - Changed from 12pt
    p_name = text_frame.add_paragraph()
    p_name.text = rfp_name_val
    p_name.font.size = Pt(18)
    p_name.font.color.rgb = RGBColor(255, 255, 255) # White text
    p_name.alignment = PP_ALIGN.LEFT

    # One blank line
    blank_p_single = text_frame.add_paragraph()
    blank_p_single.text = " "
    blank_p_single.font.size = Pt(12)


    # Today's Date (18pt, white, no key)
    p_date = text_frame.add_paragraph()
    p_date.text = current_date_val
    p_date.font.size = Pt(18)
    p_date.font.color.rgb = RGBColor(255, 255, 255) # White text
    p_date.alignment = PP_ALIGN.LEFT


    # --- Second Slide (Empty White Slide) ---
    empty_slide = prs.slides.add_slide(blank_slide_layout)
    # Apply background image, default to white for the second slide
    apply_background_image_pptx(empty_slide, second_page_bg_info["bytes"], second_page_bg_info["name"], prs, RGBColor(255, 255, 255))


    # --- Dynamic Content Slides (from 3rd slide onwards) ---
    # Define the custom blue color
    CUSTOM_BLUE_RGB = RGBColor(0x03, 0x9e, 0xed) # Hex #039eed

    # Define custom colors for keys and values in RFP Schedule
    SCHEDULE_KEY_COLOR = RGBColor(0x2E, 0x2E, 0x2E) # Dark gray (R,G,B)
    SCHEDULE_VALUE_COLOR = RGBColor(0x6A, 0x6A, 0x6A) # Lighter gray (R,G,B)

    # Define the custom table row colors
    TABLE_COLOR_PRIMARY = RGBColor(0x03, 0x9e, 0xed) # #039eed for first row (header)
    TABLE_COLOR_SECONDARY = RGBColor(0xcc, 0xdf, 0xf8) # #ccdff8 for second alternating
    TABLE_COLOR_TERTIARY = RGBColor(0xe8, 0xf0, 0xfc) # #e8f0fc for third alternating

    def set_table_row_colors(table):
        # Header row
        for cell in table.rows[0].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_COLOR_PRIMARY
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255) # White text on blue header
                    run.font.name = 'Arial Narrow' # Apply Arial Narrow to table headers

        # Data rows
        for r_idx in range(1, len(table.rows)): # Start from second row
            row = table.rows[r_idx]
            if r_idx % 2 == 1: # Odd rows (1, 3, 5...) get secondary color
                fill_color = TABLE_COLOR_SECONDARY
            else: # Even rows (2, 4, 6...) get tertiary color
                fill_color = TABLE_COLOR_TERTIARY
            
            for cell in row.cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill_color
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(0, 0, 0) # Black text on light backgrounds
                        run.font.name = 'Arial Narrow' # Apply Arial Narrow to table body text

    def add_content_to_slide(slide, title_text, content_data, is_dict=False, text_color=RGBColor(0,0,0), font_name='Arial Narrow', wrap_text=False, center_align=False): # Added center_align
        apply_background_image_pptx(slide, content_bg_info["bytes"], content_bg_info["name"], prs, RGBColor(0, 0, 0))

        title_box_content = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.7))
        title_frame_content = title_box_content.text_frame
        title_frame_content.text = title_text
        title_paragraph_content = title_frame_content.paragraphs[0] # Access first paragraph
        title_paragraph_content.font.size = Pt(24)
        title_paragraph_content.font.bold = True
        title_paragraph_content.font.color.rgb = CUSTOM_BLUE_RGB
        title_paragraph_content.font.name = font_name
        title_paragraph_content.alignment = PP_ALIGN.LEFT

        left_content = Inches(0.5)
        top_content = Inches(1.2)
        width_content = prs.slide_width - Inches(1)
        height_content = prs.slide_height - top_content - Inches(0.5)
        content_textbox = slide.shapes.add_textbox(left_content, top_content, width_content, height_content)
        content_frame = content_textbox.text_frame
        content_frame.clear()
        
        # Ensure text in content_frame wraps
        content_frame.word_wrap = True

        if is_dict and isinstance(content_data, dict):
            for key, value in content_data.items():
                p = content_frame.add_paragraph()
                p.alignment = PP_ALIGN.CENTER if center_align else PP_ALIGN.LEFT # Set alignment here

                run_key = p.add_run()  
                run_key.text = f"{key.upper()}: "  
                run_key.font.bold = True
                run_key.font.size = Pt(16)
                run_key.font.color.rgb = SCHEDULE_KEY_COLOR if title_text == "RFP Schedule" else text_color # Apply color for RFP Schedule key
                run_key.font.name = font_name

                # Special handling for 'Required Forms' and 'Evaluation Criteria' to render as bullets
                if (key == "Required Forms" or key == "Evaluation Criteria") and isinstance(value, list):
                    for item_bullet in value:
                        # For bullet points, add a new paragraph for each item
                        bullet_p = content_frame.add_paragraph()
                        bullet_p.text = str(item_bullet)
                        bullet_p.level = 1 # Indent for bullet
                        bullet_p.font.size = Pt(16)
                        bullet_p.font.color.rgb = text_color
                        bullet_p.font.name = font_name
                        bullet_p.alignment = PP_ALIGN.LEFT # Ensure bullets are left-aligned
                elif wrap_text and isinstance(value, str):
                    # Split string by newlines for better wrapping control
                    lines = value.split('\n')
                    for i, line in enumerate(lines):
                        run_value = p.add_run()  
                        run_value.text = line.strip()  
                        run_value.font.bold = False  
                        run_value.font.size = Pt(16)
                        run_value.font.color.rgb = SCHEDULE_VALUE_COLOR if title_text == "RFP Schedule" else text_color # Apply color for RFP Schedule value
                        run_value.font.name = font_name
                        # Add a new paragraph for the next line if it's not the last one
                        if i < len(lines) - 1 and lines[i+1].strip():
                            p = content_frame.add_paragraph() # Start a new paragraph for the next line for better wrapping
                            p.alignment = PP_ALIGN.CENTER if center_align else PP_ALIGN.LEFT
                else:
                    run_value = p.add_run()  
                    run_value.text = str(value)  
                    run_value.font.bold = False  
                    run_value.font.size = Pt(16)
                    run_value.font.color.rgb = SCHEDULE_VALUE_COLOR if title_text == "RFP Schedule" else text_color
                    run_value.font.name = font_name
        elif isinstance(content_data, str):
            lines_to_add = content_data.split('\n') if wrap_text else [content_data]
            for line in lines_to_add:
                if line.strip():
                    p = content_frame.add_paragraph()
                    p.text = line.strip()
                    p.font.size = Pt(16)
                    p.font.color.rgb = text_color
                    p.font.name = font_name
                    p.alignment = PP_ALIGN.CENTER if center_align else PP_ALIGN.LEFT
        else:
            p = content_frame.add_paragraph()
            p.text = str(content_data)
            p.font.size = Pt(16)
            p.font.color.rgb = text_color
            p.font.name = font_name
            p.alignment = PP_ALIGN.CENTER if center_align else PP_ALIGN.LEFT


    # --- Master List of Sections to Process and Their Order ---
    # Removed "Overall Summary", "RFP Overview Details", "Cost Proposal Form", "Scorecard Evaluation"
    NEW_SECTION_TITLES = [
        "Prospect & RFP Background",
        "Scope of Work Details",
        "Service Level Agreements (SLAs)",
        "RFP Submission Information",
        "RFP Schedule",
        "RFP Evaluation Criteria",
        "SWOT Analysis"
    ]

    # Dynamically generate slides based on the ordered list and data
    for section_title in NEW_SECTION_TITLES: # Use NEW_SECTION_TITLES directly here
        content = analysis_data.get(section_title, "No relevant content found.")
        
        if content == "No relevant content found." or (isinstance(content, (dict, list)) and not content):
            continue
        
        # We will handle slide creation for SWOT directly in its block
        if section_title != "SWOT Analysis":
            slide = prs.slides.add_slide(blank_slide_layout)  
            apply_background_image_pptx(slide, content_bg_info["bytes"], content_bg_info["name"], prs, RGBColor(0, 0, 0)) # Apply background to all content slides
        

        if section_title == "Prospect & RFP Background":
            title_box_content = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.7))
            title_frame_content = title_box_content.text_frame
            title_frame_content.text = "Prospect & RFP Background"
            title_paragraph_content = title_frame_content.paragraphs[0]
            title_paragraph_content.font.size = Pt(24)
            title_paragraph_content.font.bold = True
            title_paragraph_content.font.color.rgb = CUSTOM_BLUE_RGB
            title_paragraph_content.font.name = 'Arial Narrow'
            title_paragraph_content.alignment = PP_ALIGN.LEFT

            table_data_prospect = []
            
            client_name = analysis_data.get("Prospect & RFP Background", {}).get("Client Name", "Company")
            about_company_desc_parts = []
            bg_content = analysis_data.get("Prospect & RFP Background", {})  
            if isinstance(bg_content, dict):
                if bg_content.get('Industry'): about_company_desc_parts.append(f"Industry: {bg_content['Industry']}")
                if bg_content.get('Project Purpose'): about_company_desc_parts.append(f"Project Purpose: {bg_content['Project Purpose']}")
                if bg_content.get('Background'): about_company_desc_parts.append(f"Background: {bg_content['Background']}")
                if bg_content.get('Current Challenges'): about_company_desc_parts.append(f"Current Challenges: {bg_content['Current Challenges']}")
            elif isinstance(bg_content, str):  
                lines = [line.strip() for line in bg_content.split('\n') if line.strip()]
                for line in lines:
                    line = line.strip()
                    if line.startswith("Industry:"): about_company_desc_parts.append(line)
                    elif line.startswith("Project Purpose:"): about_company_desc_parts.append(line)
                    elif line.startswith("Background:"): about_company_desc_parts.append(line)
                    elif line.startswith("Current Challenges:"): about_company_desc_parts.append(line)
            
            about_company_desc = "\n".join(about_company_desc_parts) if about_company_desc_parts else "N/A"
            table_data_prospect.append([f"About {client_name}", about_company_desc])

            service_agreement_details_compiled = []
            
            service_agreement_raw_content = None
            if isinstance(bg_content, dict):
                service_agreement_raw_content = bg_content.get("Service Agreement Details")
            elif isinstance(bg_content, str):
                match_service_agreement = re.search(r"Service Agreement Details:\s*(.*?)(?:\n|$)", bg_content, re.IGNORECASE)
                if match_service_agreement:
                    service_agreement_raw_content = match_service_agreement.group(1).strip()
            
            if service_agreement_raw_content and service_agreement_raw_content.strip().lower() not in ['n/a', '', 'not specified', 'null', 'example.g. 5', 'example.g. 3', '5 years agreement; lbe subcontract required for 20% of work']:
                
                agreement_duration_text = []  
                years_match = re.search(r'(\d+)\s*(?:year|yr)(?:s)?', service_agreement_raw_content, re.IGNORECASE)
                months_match = re.search(r'(\d+)\s*(?:month|mo)(?:s)?', service_agreement_raw_content, re.IGNORECASE)
                
                if years_match:
                    agreement_duration_text.append(f"Agreement Years: {years_match.group(1)}")
                elif months_match:
                    num_months = int(months_match.group(1))
                    years = num_months / 12
                    agreement_duration_text.append(f"Agreement Years: {years:.1f}")  
                
                if not agreement_duration_text and ("agreement" in service_agreement_raw_content.lower() or "contract duration" in service_agreement_raw_content.lower()):
                    generic_duration_match = re.search(r'(?:(?:contract|agreement) duration|for)\s*(.*?)(?:$|;|\n)', service_agreement_raw_content, re.IGNORECASE)
                    if generic_duration_match and generic_duration_match.group(1).strip().lower() not in ['n/a', '', 'not specified', 'null']:
                        agreement_duration_text.append(f"Agreement Duration: {generic_duration_match.group(1).strip()}")
                    elif "years agreement" in service_agreement_raw_content.lower() or "months agreement" in service_agreement_raw_content.lower():
                        agreement_duration_text.append(f"Agreement Details: {service_agreement_raw_content}")


                lbe_match = re.search(r'LBE subcontract details:\s*(.*?)(?:\n|$)', service_agreement_raw_content, re.IGNORECASE)
                if lbe_match and lbe_match.group(1).strip().lower() not in ['n/a', '', 'not specified', 'null']:
                    service_agreement_details_compiled.append(f"LBE Subcontract: {lbe_match.group(1).strip()}")
                elif re.search(r'LBE subcontract', service_agreement_raw_content, re.IGNORECASE) and "LBE subcontract details" not in service_agreement_raw_content.lower():
                    service_agreement_details_compiled.append("LBE Subcontract: Mentioned, details not specified.")

                if agreement_duration_text or service_agreement_details_compiled:  
                    table_data_prospect.append(["Service Agreement", "; ".join(agreement_duration_text + service_agreement_details_compiled)])


            key_contact_data = []
            key_contact_name = None
            contact_email = None
            contact_phone = None

            if isinstance(bg_content, dict):  
                key_contact_name = bg_content.get("Key Contact Name")
                contact_email = bg_content.get("Contact Email")
                contact_phone = bg_content.get("Contact Phone")
            elif isinstance(bg_content, str):  
                match_name = re.search(r"Key Contact Name:\s*(.*?)(?:\n|$)", bg_content, re.IGNORECASE)
                if match_name: key_contact_name = match_name.group(1).strip()
                match_email = re.search(r"Contact Email:\s*(.*?)(?:\n|$)", bg_content, re.IGNORECASE)
                if match_email: contact_email = match_email.group(1).strip()
                match_phone = re.search(r"Contact Phone:\s*(.*?)(?:\n|$)", bg_content, re.IGNORECASE)
                if match_phone: contact_phone = match_phone.group(1).strip()
            
            if key_contact_name and key_contact_name.strip().lower() not in ['n/a', '', 'not specified', 'null']:  
                key_contact_data.append(key_contact_name)
            if contact_email and contact_email.strip().lower() not in ['n/a', '', 'not specified', 'null', 'example.com/email', 'john.doe@example.com']:  
                key_contact_data.append(f"Email: {contact_email}")
            if contact_phone and contact_phone.strip().lower() not in ['n/a', '', 'not specified', 'null', '555-123-4567']:  
                key_contact_data.append(f"Phone: {contact_phone}")

            if key_contact_data:  
                table_data_prospect.append(["Key Contact", "\n".join(key_contact_data)])


            virtual_meeting_output = []
            if isinstance(bg_content, dict):  
                vm_details = bg_content.get("Virtual Meeting Details")
                if isinstance(vm_details, dict) and vm_details:  
                    for key_in_vm in ['Meeting Topic', 'Meeting ID', 'Passcode', 'Join Link']:
                        value_in_vm = vm_details.get(key_in_vm)
                        if value_in_vm and str(value_in_vm).strip().lower() not in ['n/a', '', 'not specified', 'null', 'example.com/meeting', 'https://example.com/meeting', 'example.com/zoom', 'https://actual-meeting-link.com/join', '123-456-7890', 'rfp2025']:
                            display_key = "Link" if key_in_vm == "Join Link" else key_in_vm.replace('Meeting ', '')
                            virtual_meeting_output.append(f"{display_key}: {value_in_vm}")  
                elif isinstance(vm_details, str) and vm_details.strip().lower() not in ['n/a', '', 'not specified', 'null', 'https://example.com/meeting', 'example.com/zoom', 'https://actual-meeting-link.com/join']:
                    virtual_meeting_output.append(vm_details)
            
            if virtual_meeting_output:  
                table_data_prospect.append(["Virtual Meeting Info", "\n".join(virtual_meeting_output)])


            rows_prospect = len(table_data_prospect) + 1  
            cols_prospect = 2
            left_prospect = Inches(0.5)
            top_prospect = Inches(1.5)
            width_prospect = prs.slide_width - Inches(1)
            height_prospect = prs.slide_height - top_prospect - Inches(0.5)

            if len(table_data_prospect) > 0:
                table = slide.shapes.add_table(rows_prospect, cols_prospect, left_prospect, top_prospect, width_prospect, height_prospect).table

                table.columns[0].width = Inches(2.5)  
                table.columns[1].width = width_prospect - Inches(2.5)  

                hdr_cells = table.rows[0].cells
                hdr_cells[0].text = "#"
                hdr_cells[1].text = "Description"
                for c_idx in range(cols_prospect):
                    cell = hdr_cells[c_idx]
                    text_frame = cell.text_frame
                    p = text_frame.paragraphs[0]
                    p.font.bold = True
                    p.font.size = Pt(14)
                    p.font.name = 'Arial Narrow'
                    p.alignment = PP_ALIGN.LEFT

                for r_idx, row_data in enumerate(table_data_prospect):
                    for c_idx, cell_data in enumerate(row_data):
                        cell = table.cell(r_idx + 1, c_idx)
                        text_frame = cell.text_frame
                        text_frame.clear()
                        p = text_frame.paragraphs[0]
                        p.text = str(cell_data)
                        p.font.size = Pt(12)
                        p.font.name = 'Arial Narrow'
                        p.alignment = PP_ALIGN.LEFT
                
                set_table_row_colors(table)
            continue  

        if section_title == "Scope of Work Details":
            scope_raw_content = analysis_data["Scope of Work Details"]
            if scope_raw_content != "No relevant content found." and (isinstance(scope_raw_content, dict) or (isinstance(scope_raw_content, str) and scope_raw_content.strip())):
                
                scope_parsed_data = []

                if isinstance(scope_raw_content, dict):
                    for area, desc in scope_raw_content.items():
                        formatted_desc = ""
                        if isinstance(desc, list):
                            formatted_desc = "\n".join(desc)
                        elif isinstance(desc, dict):
                            formatted_desc = json.dumps(desc, indent=2)
                        else:
                            formatted_desc = str(desc)
                        scope_parsed_data.append([area, formatted_desc])
                elif isinstance(scope_raw_content, str):
                    lines = [line.strip() for line in scope_raw_content.split('\n') if line.strip()]
                    current_area = ""
                    current_description = []
                    for line in lines:
                        if line.endswith(":") or line.endswith("—") or (len(line) > 0 and line[0].isupper() and (len(line.split()) < 5 or line.endswith('.'))): # Heuristic for a new 'Area'
                            if current_area and current_description:
                                scope_parsed_data.append([current_area, "\n".join(current_description)])
                            current_area = re.sub(r"[:—]$", "", line).strip()
                            current_description = []
                        else:
                            current_description.append(line)
                    if current_area and current_description: # Add the last one
                        scope_parsed_data.append([current_area, "\n".join(current_description)])

                if not scope_parsed_data:
                    scope_parsed_data.append(["No specific areas found", "No detailed description extracted. Please refer to the full RFP document."])

                max_rows_per_slide = 8  
                scope_chunks_for_slides = [scope_parsed_data[i:i + max_rows_per_slide] for i in range(0, len(scope_parsed_data), max_rows_per_slide)]

                for i, chunk in enumerate(scope_chunks_for_slides):
                    if i > 0:
                        slide = prs.slides.add_slide(blank_slide_layout)  
                        apply_background_image_pptx(slide, content_bg_info["bytes"], content_bg_info["name"], prs, RGBColor(0, 0, 0)) # Re-apply background for new chunk slides
                    
                    title_text = "Scope of Work"
                    if i > 0:  
                        title_text += " - Continued"

                    title_box_content = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.7))
                    title_frame_content = title_box_content.text_frame
                    title_frame_content.text = title_text
                    title_paragraph_content = title_frame_content.paragraphs[0]
                    title_paragraph_content.font.size = Pt(24)
                    title_paragraph_content.font.bold = True
                    title_paragraph_content.font.color.rgb = CUSTOM_BLUE_RGB
                    title_paragraph_content.font.name = 'Arial Narrow'
                    title_paragraph_content.alignment = PP_ALIGN.LEFT


                    rows_current_chunk = len(chunk) + 1
                    cols = 2
                    left = Inches(0.5)
                    top = Inches(1.5)
                    width = Inches(9)
                    height = Inches(5.5)

                    table = slide.shapes.add_table(rows_current_chunk, cols, left, top, width, height).table

                    table.columns[0].width = Inches(2.5)
                    table.columns[1].width = Inches(6.5)

                    hdr_cells = table.rows[0].cells
                    hdr_cells[0].text = "Areas"
                    hdr_cells[1].text = "Description"
                    for c_idx in range(cols):
                        cell = hdr_cells[c_idx]
                        text_frame = cell.text_frame
                        p = text_frame.paragraphs[0]
                        p.font.bold = True
                        p.font.size = Pt(12)
                        p.font.name = 'Arial Narrow'
                        p.alignment = PP_ALIGN.LEFT

                    for r_idx, row_data in enumerate(chunk):
                        for c_idx, cell_data in enumerate(row_data):
                            cell = table.cell(r_idx + 1, c_idx)
                            text_frame = cell.text_frame
                            text_frame.clear()
                            p = text_frame.paragraphs[0]
                            p.text = str(cell_data)
                            p.font.size = Pt(12)
                            p.font.name = 'Arial Narrow'
                            p.alignment = PP_ALIGN.LEFT
                    
                    set_table_row_colors(table)
            continue  

        sla_content_data = analysis_data.get("Service Level Agreements (SLAs)", None)  
        
        has_meaningful_sla_content = False
        if isinstance(sla_content_data, list) and len(sla_content_data) > 0:
            for item in sla_content_data:
                if isinstance(item, dict) and any(value and str(value).strip().lower() not in ['n/a', '', 'not specified', 'null'] for value in item.values()):
                    has_meaningful_sla_content = True
                    break
        elif isinstance(sla_content_data, dict) and sla_content_data:
            if any(value and str(value).strip().lower() not in ['n/a', '', 'not specified', 'null'] for value in sla_content_data.values()):
                has_meaningful_sla_content = True
        elif isinstance(sla_content_data, str) and sla_content_data.strip().lower() not in ['n/a', '', 'not specified', 'null']:
            has_meaningful_sla_content = True

        if section_title == "Service Level Agreements (SLAs)" and has_meaningful_sla_content:
            sla_details_for_table = []
            client_current_key = "Client's Current"
            vendor_responsibility_key = "Vendor's Responsibility"

            if isinstance(sla_content_data, list):
                for item in sla_content_data:
                    if isinstance(item, dict):
                        sla_line = []
                        if item.get('Metric') and item['Metric'].strip().lower() not in ['n/a', '', 'not specified', 'null']:
                            sla_line.append(f"Metric: {item['Metric']}")
                        if item.get(client_current_key) and item.get(client_current_key).strip().lower() not in ['n/a', '', 'not specified', 'null']:
                            sla_line.append(f"Client's Current: {item.get(client_current_key)}")
                        if item.get(vendor_responsibility_key) and item.get(vendor_responsibility_key).strip().lower() not in ['n/a', '', 'not specified', 'null']:
                            sla_line.append(f"Vendor's Resp.: {item.get(vendor_responsibility_key)}")
                        
                        if sla_line:
                            sla_details_for_table.append("; ".join(sla_line))
            elif isinstance(sla_content_data, dict):
                sla_line = []
                if sla_content_data.get('Metric') and sla_content_data['Metric'].strip().lower() not in ['n/a', '', 'not specified', 'null']:
                    sla_line.append(f"Metric: {sla_content_data['Metric']}")
                if sla_content_data.get(client_current_key) and sla_content_data.get(client_current_key).strip().lower() not in ['n/a', '', 'not specified', 'null']:
                    sla_line.append(f"Client's Current: {sla_content_data.get(client_current_key)}")
                if sla_content_data.get(vendor_responsibility_key) and sla_content_data.get(vendor_responsibility_key).strip().lower() not in ['n/a', '', 'not specified', 'null']:
                    sla_line.append(f"Vendor's Resp.: {sla_content_data.get(vendor_responsibility_key)}")
                
                if sla_line:
                    sla_details_for_table.append("; ".join(sla_line))
            elif isinstance(sla_content_data, str):
                sla_details_for_table.append(str(sla_content_data).strip())

            if sla_details_for_table:
                # Reuse the existing slide as it's already created and background applied
                add_content_to_slide(slide, "Service Level Agreements (SLAs)", "\n".join(sla_details_for_table), wrap_text=True)
            continue  

    # --- RFP Schedule (Centralized with specific colors) ---
        if section_title == "RFP Schedule":
            add_content_to_slide(slide, section_title, content, is_dict=isinstance(content, dict), center_align=True, wrap_text=True) # Added wrap_text
            continue  

    # --- Standard Content Slides (handled by generic add_content_to_slide) ---
        if section_title == "RFP Submission Information":
            submission_info_content = analysis_data.get("RFP Submission Information", {})
            if isinstance(submission_info_content, dict):
                title_box_content = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.7))
                title_frame_content = title_box_content.text_frame
                title_frame_content.text = section_title
                title_paragraph_content = title_frame_content.paragraphs[0]
                title_paragraph_content.font.size = Pt(24)
                title_paragraph_content.font.bold = True
                title_paragraph_content.font.color.rgb = CUSTOM_BLUE_RGB
                title_paragraph_content.font.name = 'Arial Narrow'
                title_paragraph_content.alignment = PP_ALIGN.LEFT

                left_content = Inches(0.5)
                top_content = Inches(1.2)
                width_content = prs.slide_width - Inches(1)
                height_content = prs.slide_height - top_content - Inches(0.5)
                content_textbox = slide.shapes.add_textbox(left_content, top_content, width_content, height_content)
                content_frame = content_textbox.text_frame
                content_frame.clear()
                content_frame.word_wrap = True # Ensure wrapping for submission info

                for key, value in submission_info_content.items():
                    p = content_frame.add_paragraph()
                    run_key = p.add_run()  
                    run_key.text = f"{key.upper()}: "  
                    run_key.font.bold = True
                    run_key.font.size = Pt(16)
                    run_key.font.color.rgb = RGBColor(0, 0, 0)
                    run_key.font.name = 'Arial Narrow'

                    if key == "Required Forms" and isinstance(value, list):
                        for form_item in value:
                            bullet_p = content_frame.add_paragraph()
                            bullet_p.text = str(form_item)
                            bullet_p.level = 1  
                            bullet_p.font.size = Pt(16)
                            bullet_p.font.color.rgb = RGBColor(0, 0, 0)
                            bullet_p.font.name = 'Arial Narrow'
                    elif isinstance(value, str):
                        # Split by newline for better wrapping
                        lines = value.split('\n')
                        for i, line in enumerate(lines):
                            run_value = p.add_run()  
                            run_value.text = line.strip()  
                            run_value.font.bold = False
                            run_value.font.size = Pt(16)
                            run_value.font.color.rgb = RGBColor(0, 0, 0)
                            run_value.font.name = 'Arial Narrow'
                            if i < len(lines) - 1 and lines[i+1].strip():
                                p = content_frame.add_paragraph() # New paragraph for next line
                    else:
                        run_value = p.add_run()  
                        run_value.text = str(value)  
                        run_value.font.bold = False
                        run_value.font.size = Pt(16)
                        run_value.font.color.rgb = RGBColor(0, 0, 0)
                        run_value.font.name = 'Arial Narrow'
            continue  

        elif section_title == "RFP Evaluation Criteria":
            evaluation_criteria_content = analysis_data.get("RFP Evaluation Criteria", {})
            if isinstance(evaluation_criteria_content, dict):
                title_box_content = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.7))
                title_frame_content = title_box_content.text_frame
                title_frame_content.text = section_title
                title_paragraph_content = title_frame_content.paragraphs[0]
                title_paragraph_content.font.size = Pt(24)
                title_paragraph_content.font.bold = True
                title_paragraph_content.font.color.rgb = CUSTOM_BLUE_RGB
                title_paragraph_content.font.name = 'Arial Narrow'
                title_paragraph_content.alignment = PP_ALIGN.LEFT

                left_content = Inches(0.5)
                top_content = Inches(1.2)
                width_content = prs.slide_width - Inches(1)
                height_content = prs.slide_height - top_content - Inches(0.5)
                content_textbox = slide.shapes.add_textbox(left_content, top_content, width_content, height_content)
                content_frame = content_textbox.text_frame
                content_frame.clear()
                content_frame.word_wrap = True # Ensure wrapping for evaluation criteria

                for key, value in evaluation_criteria_content.items():
                    p = content_frame.add_paragraph()
                    run_key = p.add_run()
                    run_key.text = f"{key.upper()}: "
                    run_key.font.bold = True
                    run_key.font.size = Pt(16)
                    run_key.font.color.rgb = RGBColor(0, 0, 0)
                    run_key.font.name = 'Arial Narrow'

                    if key == "Evaluation Criteria" and isinstance(value, list):
                        for criterion_item in value:
                            bullet_p = content_frame.add_paragraph()
                            bullet_p.text = str(criterion_item)
                            bullet_p.level = 1
                            bullet_p.font.size = Pt(16)
                            bullet_p.font.color.rgb = RGBColor(0, 0, 0)
                            bullet_p.font.name = 'Arial Narrow'
                    elif isinstance(value, str):
                        lines = value.split('\n')
                        for i, line in enumerate(lines):
                            run_value = p.add_run()
                            run_value.text = line.strip()
                            run_value.font.bold = False
                            run_value.font.size = Pt(16)
                            run_value.font.color.rgb = RGBColor(0, 0, 0)
                            run_value.font.name = 'Arial Narrow'
                            if i < len(lines) - 1 and lines[i+1].strip():
                                p = content_frame.add_paragraph() # New paragraph for next line
                    else:
                        run_value = p.add_run()
                        run_value.text = str(value)
                        run_value.font.bold = False
                        run_value.font.size = Pt(16)
                        run_value.font.color.rgb = RGBColor(0, 0, 0)
                        run_value.font.name = 'Arial Narrow'
            continue

        elif section_title == "SWOT Analysis":
            # Start of your boilerplate code for SWOT
            slide_layout = prs.slide_layouts[1] # Assuming layout 1 is a title and content layout
            slide = prs.slides.add_slide(slide_layout)
            apply_background_image_pptx(slide, content_bg_info["bytes"], content_bg_info["name"], prs, RGBColor(0, 0, 0)) # Apply background
            
            # --- Explicitly adding the title ---
            # Use a separate text box for the title for better control, matching other slides' titles
            title_box_custom = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.7))
            title_frame_custom = title_box_custom.text_frame
            title_frame_custom.text = "SWOT Analysis"
            title_paragraph_custom = title_frame_custom.paragraphs[0]
            title_paragraph_custom.font.size = Pt(24)
            title_paragraph_custom.font.bold = True
            title_paragraph_custom.font.color.rgb = CUSTOM_BLUE_RGB
            title_paragraph_custom.font.name = 'Arial Narrow'
            title_paragraph_custom.alignment = PP_ALIGN.LEFT
            # --- End title addition ---

            # Clear existing content placeholder if layout 1 has one (usually it does)
            for shape_placeholder in slide.placeholders:
                if shape_placeholder.is_placeholder and shape_placeholder.has_text_frame:
                    # Often the main content placeholder is idx 1
                    if shape_placeholder.placeholder_format.idx == 1:
                        shape_placeholder.text_frame.clear()

            swot_raw_content = analysis_data.get("SWOT Analysis", {})  

            if isinstance(swot_raw_content, dict):
                strengths = swot_raw_content.get('Strengths', [])
                weaknesses = swot_raw_content.get('Weaknesses', [])
                opportunities = swot_raw_content.get('Opportunities', [])
                threats = swot_raw_content.get('Threats', [])
                
                # Check if lists contain actual content
                if any(strengths) or any(weaknesses) or any(opportunities) or any(threats):
                    # Use Inches() for dimensions
                    left = Inches(1.0)
                    top = Inches(1.5)
                    width = Inches(4)
                    height = Inches(2.5)

                    # Ensure all lists are actually lists, convert if needed
                    strengths = [str(item) for item in (strengths if isinstance(strengths, list) else [strengths])]
                    weaknesses = [str(item) for item in (weaknesses if isinstance(weaknesses, list) else [weaknesses])]
                    opportunities = [str(item) for item in (opportunities if isinstance(opportunities, list) else [opportunities])]
                    threats = [str(item) for item in (threats if isinstance(threats, list) else [threats])]

                    # Helper function for adding a quadrant with proper bullets
                    def add_swot_quadrant(slide, x, y, w, h, title, items):
                        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(240, 240, 240) # Light grey background
                        shape.line.fill.background() # No border line by default, or set a color if desired

                        text_frame = shape.text_frame
                        text_frame.clear()
                        text_frame.word_wrap = True # Ensure text wraps within the shape

                        # Add title for the quadrant
                        p_title = text_frame.add_paragraph()
                        run_title = p_title.add_run()
                        run_title.text = f"{title}:"
                        run_title.font.size = Pt(16)
                        run_title.font.bold = True
                        run_title.font.color.rgb = CUSTOM_BLUE_RGB # Use your custom blue for quadrant titles
                        p_title.alignment = PP_ALIGN.CENTER
                        
                        # Add a small spacer line after title for better visual separation
                        p_spacer = text_frame.add_paragraph()
                        p_spacer.text = ""
                        p_spacer.font.size = Pt(6) # Smaller font for minimal spacing

                        # Add bullet points for content
                        for item in items:
                            if item.strip():
                                p_bullet = text_frame.add_paragraph()
                                p_bullet.text = item.strip()
                                p_bullet.level = 1 # This makes it a bullet point
                                p_bullet.font.size = Pt(12) # Slightly smaller font for list items
                                p_bullet.font.color.rgb = RGBColor(0, 0, 0) # Black text
                                p_bullet.alignment = PP_ALIGN.LEFT


                    # Top-Left: Strengths
                    add_swot_quadrant(slide, left, top, width, height, "Strengths", strengths)

                    # Top-Right: Weaknesses
                    add_swot_quadrant(slide, left + width + Inches(0.5), top, width, height, "Weaknesses", weaknesses)

                    # Bottom-Left: Opportunities
                    add_swot_quadrant(slide, left, top + height + Inches(0.5), width, height, "Opportunities", opportunities)

                    # Bottom-Right: Threats
                    add_swot_quadrant(slide, left + width + Inches(0.5), top + height + Inches(0.5), width, height, "Threats", threats)

                else: # Fallback if no SWOT content is found
                    # Use a general content textbox if no structured SWOT data
                    content_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), prs.slide_width - Inches(1), prs.slide_height - Inches(2))
                    content_text_frame = content_text_box.text_frame
                    content_text_frame.word_wrap = True
                    p = content_text_frame.add_paragraph()
                    p.text = "SWOT Analysis content from RFP (no specific categories found):\n" + json.dumps(swot_raw_content, indent=2)
                    p.font.size = Pt(16)
                    p.alignment = PP_ALIGN.LEFT
            else: # Fallback if SWOT content is not a dict
                # Use a general content textbox if no structured SWOT data
                content_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), prs.slide_width - Inches(1), prs.slide_height - Inches(2))
                content_text_frame = content_text_box.text_frame
                content_text_frame.word_wrap = True
                p = content_text_frame.add_paragraph()
                p.text = "SWOT Analysis content from RFP:\n" + str(swot_raw_content).strip()
                p.font.size = Pt(16)
                p.alignment = PP_ALIGN.LEFT
            continue # Ensure we don't try to process this as a generic section after custom handling


    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

def generate_docx_option1(analysis_data):
    """Generates a DOCX file from the analysis results based on new sections"""
    doc = docx.Document()

    # --- Table of Contents ---
    doc.add_heading("Table of Contents", level=1)
    for i, title in enumerate(NEW_SECTION_TITLES): # Uses updated NEW_SECTION_TITLES
        p = doc.add_paragraph()
        p.add_run(f"{i+1}. {title}").bold = True
        p.paragraph_format.left_indent = DocxInches(0.5)
    doc.add_page_break()


    # --- Prospect & RFP Background (Special Handling for Dictionary Output) ---
    title_background = "Prospect & RFP Background"
    if title_background in analysis_data:
        content_background = analysis_data[title_background]
        doc.add_heading(title_background, level=1)
        
        if isinstance(content_background, dict):
            for key, value in content_background.items():
                p = doc.add_paragraph()
                run_key = p.add_run(f"{key.upper()}: ")
                run_key.bold = True
                p.add_run(str(value))  
        elif isinstance(content_background, str):
            for line in content_background.split('\n'):
                if line.strip():
                    doc.add_paragraph(line.strip())
        else:
            doc.add_paragraph(str(content_background))
        doc.add_page_break()


    # --- Scope of Work (Table) ---
    doc.add_heading("Scope of Work Details", level=1)
    table_scope = doc.add_table(rows=1, cols=2)
    table_scope.style = 'Table Grid'
    hdr_cells = table_scope.rows[0].cells
    hdr_cells[0].text = "Areas"
    hdr_cells[1].text = "Description"

    for cell in hdr_cells:
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    scope_raw_content = analysis_data.get("Scope of Work Details", "")
    scope_parsed_data = []

    if isinstance(scope_raw_content, dict):
        for area, desc in scope_raw_content.items():
            formatted_desc = ""
            if isinstance(desc, list):
                formatted_desc = "\n".join(desc)
            elif isinstance(desc, dict):
                formatted_desc = json.dumps(desc, indent=2)
            else:
                formatted_desc = str(desc)
            scope_parsed_data.append([area, formatted_desc])
    elif isinstance(scope_raw_content, str):
        lines = [line.strip() for line in scope_raw_content.split('\n') if line.strip()]
        current_area = ""
        current_description = []
        for line in lines:
            if line.endswith(":") or line.endswith("—") or (len(line) > 0 and line[0].isupper() and (len(line.split()) < 5 or line.endswith('.'))): # Heuristic for a new 'Area'
                if current_area and current_description:
                    scope_parsed_data.append([current_area, "\n".join(current_description)])
                current_area = re.sub(r"[:—]$", "", line).strip()
                current_description = []
            else:
                current_description.append(line)
        if current_area and current_description: # Add the last one
            scope_parsed_data.append([current_area, "\n".join(current_description)])
    
    if not scope_parsed_data:
        scope_parsed_data.append(["No specific areas found", "No detailed description extracted. Please refer to the full RFP document."])

    for area, description in scope_parsed_data:
        row_cells = table_scope.add_row().cells
        row_cells[0].text = str(area)
        row_cells[1].text = str(description)

    doc.add_page_break()

    # --- Service Level Agreements (SLAs) ---
    doc.add_heading("Service Level Agreements (SLAs)", level=1)
    sla_content_data = analysis_data.get("Service Level Agreements (SLAs)", "No relevant content found.")
    
    if isinstance(sla_content_data, list):
        if sla_content_data:
            table_sla = doc.add_table(rows=1, cols=3)
            table_sla.style = 'Table Grid'
            hdr_cells_sla = table_sla.rows[0].cells
            hdr_cells_sla[0].text = "Metric"
            hdr_cells_sla[1].text = "Client's Current"
            hdr_cells_sla[2].text = "Vendor's Responsibility"

            for cell in hdr_cells_sla:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.bold = True
                    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

            for item in sla_content_data:
                if isinstance(item, dict):
                    row_cells = table_sla.add_row().cells
                    row_cells[0].text = item.get('Metric', 'N/A')
                    row_cells[1].text = item.get("Client's Current", 'N/A')  
                    row_cells[2].text = item.get("Vendor's Responsibility", 'N/A')
        else:
            doc.add_paragraph("No relevant content found for Service Level Agreements (SLAs).")
    elif isinstance(sla_content_data, str) and sla_content_data.strip():
        doc.add_paragraph(sla_content_data)
    else:
        doc.add_paragraph("No relevant content found for Service Level Agreements (SLAs).")
    doc.add_page_break()


    # --- Dynamic Sections (Formatted as KEY: VALUE or bullet points) ---
    sections_to_format_as_key_value = [
        "RFP Submission Information",
        "RFP Schedule",
        "RFP Evaluation Criteria"
    ]

    for title in sections_to_format_as_key_value:
        content_from_analysis = analysis_data.get(title, "No relevant content found.")
        if content_from_analysis != "No relevant content found." and \
           ((isinstance(content_from_analysis, dict) and content_from_analysis) or \
            (isinstance(content_from_analysis, str) and content_from_analysis.strip())):
            
            doc.add_heading(title, level=1)
            if isinstance(content_from_analysis, dict):
                for key, value in content_from_analysis.items():
                    p = doc.add_paragraph()
                    run_key = p.add_run(f"{key.upper()}: ")
                    run_key.bold = True
                    if isinstance(value, list):
                        for item in value:
                            if str(item).strip():
                                doc.add_paragraph(str(item).strip(), style='List Bullet')
                    else:
                        p.add_run(str(value))
            else: 
                content = str(content_from_analysis)
                if "\n" in content:
                    for line in content.split('\n'):
                        if line.strip():
                            p = doc.add_paragraph(line.strip(), style='List Bullet')
                else:
                    doc.add_paragraph(content)
            doc.add_page_break()


    # --- SWOT Analysis (Table for DOCX) ---
    doc.add_heading("SWOT Analysis", level=1)
    swot_raw_content = analysis_data.get("SWOT Analysis", {})

    if isinstance(swot_raw_content, dict) and (swot_raw_content.get('Strengths') or swot_raw_content.get('Weaknesses') or
                                               swot_raw_content.get('Opportunities') or swot_raw_content.get('Threats')):
        table_swot = doc.add_table(rows=2, cols=2)
        table_swot.style = 'Table Grid'

        hdr_s = table_swot.cell(0, 0).paragraphs[0]
        hdr_s.add_run("Strengths").bold = True
        hdr_s.alignment = WD_ALIGN_PARAGRAPH.CENTER

        hdr_w = table_swot.cell(0, 1).paragraphs[0]
        hdr_w.add_run("Weaknesses").bold = True
        hdr_w.alignment = WD_ALIGN_PARAGRAPH.CENTER

        hdr_o = table_swot.cell(1, 0).paragraphs[0]
        hdr_o.add_run("Opportunities").bold = True
        hdr_o.alignment = WD_ALIGN_PARAGRAPH.CENTER

        hdr_t = table_swot.cell(1, 1).paragraphs[0]
        hdr_t.add_run("Threats").bold = True
        hdr_t.alignment = WD_ALIGN_PARAGRAPH.CENTER

        def add_bullets_to_docx_cell(cell, items):
            if isinstance(items, list):
                for item in items:
                    if str(item).strip():
                        cell.add_paragraph(str(item).strip(), style='List Bullet')
            elif str(items).strip():
                cell.add_paragraph(str(items).strip(), style='List Bullet')

        add_bullets_to_docx_cell(table_swot.cell(0, 0), swot_raw_content.get('Strengths', []))
        add_bullets_to_docx_cell(table_swot.cell(0, 1), swot_raw_content.get('Weaknesses', []))
        add_bullets_to_docx_cell(table_swot.cell(1, 0), swot_raw_content.get('Opportunities', []))
        add_bullets_to_docx_cell(table_swot.cell(1, 1), swot_raw_content.get('Threats', []))
    else:
        doc.add_paragraph("No detailed SWOT analysis found.")
    doc.add_page_break()


    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf

# --- NEW FUNCTIONS FOR "UNDERSTANDING DOC" ---

def analyze_for_understanding_doc(context_text, model_name_for_analysis="deepseek-r1-distill-llama-70b"):
    """
    Analyzes the document to extract content for the "Understanding Doc" sections,
    returning a structured JSON object (Python dictionary).
    Enhanced to request more detail and specific bullet point counts.
    """
    client = Groq(api_key=os.getenv("GROQ_API_KEY"))

    prompt = f"""
    Based on the following RFP document content, extract and summarize the information into a structured JSON format.
    Provide comprehensive and detailed information for each section as specified.

    Document Content:
    {context_text[:8000]}

    Desired JSON Output Format:
    {{
      "Understanding Objectives": {{
        "Strategic Objectives": [
          "Provide at least 20 detailed bullet points outlining the strategic objectives from the RFP. Each point should be complete, actionable, and clearly tied to business value or project goals."
        ],
        "LLM Understanding": [
          "Provide at least 15–20 detailed bullet points about how an LLM would interpret and understand this RFP, including inferred requirements, challenges, and architectural complexity."
        ],
        "Company Benefits": [
          "Provide at least 20 bullet points explaining benefits to the company if this project is implemented successfully. Address strategic, operational, financial, technological, and competitive advantages."
        ]
      }},

      "Functionality Mapping": {{
        "description": "Generate a report-style write-up with clearly labeled subsections for each required functionality. Each section must be a minimum of 75–100 words, using paragraphs and bullets where helpful. Use professional, clear, and formal language suitable for both technical and business stakeholders.",
        "functionalities": [
          {{
            "name": "Functionality Name 1",
            "purpose": "Full sentence explanation of the business or technical objective.",
            "solution": "Description of the approach, process, or module proposed to meet the need.",
            "technical_components": "Tools, platforms, infrastructure, APIs, or frameworks involved.",
            "dependencies_constraints": "Business rules, timeline constraints, or integration needs.",
            "expected_outcomes": "Operational benefits, KPIs, or impact of implementing this functionality."
          }},
          {{
            "name": "Functionality Name 2",
            "purpose": "",
            "solution": "",
            "technical_components": "",
            "dependencies_constraints": "",
            "expected_outcomes": ""
          }}
        ]
      }},

      "Architecture Design": {{
        "description": "Provide an architecture overview broken down by major system components or modules. Each component section should be 75–100+ words.",
        "components": [
          {{
            "name": "Component Name 1 (e.g., Backend Services, Data Layer, API Gateway)",
            "function": "Role of the component in the system.",
            "technologies_rationale": "Tools, platforms, or frameworks proposed and justification.",
            "data_flow": "How data flows through this component (input -> processing -> output).",
            "integration_points": "External/internal systems, APIs, protocols, or middleware.",
            "scalability_redundancy": "Horizontal/vertical scaling, load balancing, failover design.",
            "security": "Authentication, encryption, compliance (e.g., GDPR, HIPAA).",
            "performance_optimization": "Caching, indexing, API throttling, etc.",
            "future_proofing": "Modularity, versioning, extensibility support."
          }},
          {{
            "name": "Component Name 2 (e.g., Data Layer)",
            "function": "",
            "technologies_rationale": "",
            "data_flow": "",
            "integration_points": "",
            "scalability_redundancy": "",
            "security": "",
            "performance_optimization": "",
            "future_proofing": ""
          }}
        ]
      }},

      "Work Plan (Detail)": {{
        "description": "Break the project into clear phases and document each in depth. Each phase must include at least 75–100 words total.",
        "phases": [
          {{
            "name": "Discovery",
            "tasks_subtasks": ["Bullet points with short descriptions."],
            "effort_estimate": "4 weeks (Person-days or weeks).",
            "roles_responsibilities": "Who is doing what.",
            "dependencies_prerequisites": "What must be ready before this phase.",
            "milestones_deliverables": "Key outputs.",
            "tools_methodologies": "Agile, Jira, GitHub, CI/CD.",
            "risk_mitigation": "Fallback plans, assumptions.",
            "client_involvement": "Review points, feedback gates."
          }},
          {{
            "name": "Design",
            "tasks_subtasks": [],
            "effort_estimate": "",
            "roles_responsibilities": "",
            "dependencies_prerequisites": "",
            "milestones_deliverables": "",
            "tools_methodologies": "",
            "risk_mitigation": "",
            "client_involvement": ""
          }},
          {{
            "name": "Development",
            "tasks_subtasks": [],
            "effort_estimate": "",
            "roles_responsibilities": "",
            "dependencies_prerequisites": "",
            "milestones_deliverables": "",
            "tools_methodologies": "",
            "risk_mitigation": "",
            "client_involvement": ""
          }},
          {{
            "name": "Testing",
            "tasks_subtasks": [],
            "effort_estimate": "",
            "roles_responsibilities": "",
            "dependencies_prerequisites": "",
            "milestones_deliverables": "",
            "tools_methodologies": "",
            "risk_mitigation": "",
            "client_involvement": ""
          }},
          {{
            "name": "Deployment",
            "tasks_subtasks": [],
            "effort_estimate": "",
            "roles_responsibilities": "",
            "dependencies_prerequisites": "",
            "milestones_deliverables": "",
            "tools_methodologies": "",
            "risk_mitigation": "",
            "client_involvement": ""
          }},
          {{
            "name": "Maintenance & Support",
            "tasks_subtasks": [],
            "effort_estimate": "",
            "roles_responsibilities": "",
            "dependencies_prerequisites": "",
            "milestones_deliverables": "",
            "tools_methodologies": "",
            "risk_mitigation": "",
            "client_involvement": ""
          }}
        ]
      }},

      "Cost & Commercials": {{
        "description": "Present a detailed cost plan broken down by category. Include both narrative explanations and itemized cost details for each category, and a final summary table.",
        "categories": [
          {{
            "name": "Infrastructure Setup",
            "narrative": "Detailed narrative with assumptions (e.g., '1 DevOps Engineer @ $100/hr for 6 months').",
            "cost_per_unit": "500000",
            "fixed_vs_variable": "Fixed",
            "payment_milestones": ["30% on kickoff", "40% after UAT"],
            "scalability_factors": "Per-seat cost growth."
          }},
          {{
            "name": "Technology Stack Licensing",
            "narrative": "",
            "cost_per_unit": "",
            "fixed_vs_variable": "",
            "payment_milestones": [],
            "scalability_factors": ""
          }},
          {{
            "name": "Human Resources (by role/seniority)",
            "narrative": "",
            "cost_per_unit": "",
            "fixed_vs_variable": "",
            "payment_milestones": [],
            "scalability_factors": ""
          }},
          {{
            "name": "Software/Tools Subscription",
            "narrative": "",
            "cost_per_unit": "",
            "fixed_vs_variable": "",
            "payment_milestones": [],
            "scalability_factors": ""
          }},
          {{
            "name": "Operations & Maintenance",
            "narrative": "",
            "cost_per_unit": "",
            "fixed_vs_variable": "",
            "payment_milestones": [],
            "scalability_factors": ""
          }},
          {{
            "name": "Contingency Buffer",
            "narrative": "",
            "cost_per_unit": "",
            "fixed_vs_variable": "",
            "payment_milestones": [],
            "scalability_factors": ""
          }}
        ],
        "total_cost_summary_table": [
          {{"category": "Infrastructure Setup", "estimated_cost_usd": "500000"}},
          {{"category": "Technology Stack Licensing", "estimated_cost_usd": ""}},
          {{"category": "Human Resources", "estimated_cost_usd": ""}},
          {{"category": "Software/Tools Subscription", "estimated_cost_usd": ""}},
          {{"category": "Operations & Maintenance", "estimated_cost_usd": ""}},
          {{"category": "Contingency Buffer", "estimated_cost_usd": ""}},
          {{"category": "Overall Total", "estimated_cost_usd": ""}}
        ],
        "currency": "USD"
      }}
    }}
    """

    try:
        response = client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model=model_name_for_analysis,
            response_format={"type": "json_object"},
            temperature=0.3
        )
        return json.loads(response.choices[0].message.content)
    except json.JSONDecodeError as e:
        st.error(f"Error decoding JSON response from LLM for Understanding Doc: {e}")
        st.error(f"Raw response: {response.choices[0].message.content}")
        return {}
    except Exception as e:
        st.error(f"An unexpected error occurred during LLM analysis for Understanding Doc: {str(e)}")
        if "GROQ_API_KEY" not in os.environ:
            st.error("GROQ_API_KEY environment variable not found. Please set it in your .env file or Streamlit secrets.")
        return {}

    
# --- TECHNICAL PROPOSAL DOCUMENT GENERATION FUNCTIONS ---

def create_main_header(section, project_name):
    """Creates the main header for the document."""
    header = section.header
    header.is_linked_to_previous = False
    table = header.add_table(rows=1, cols=2, width=Inches(6.5))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    
    # Left cell for text
    left_cell = table.cell(0, 0)
    left_cell.width = Inches(4.5)
    p_left = left_cell.paragraphs[0]
    p_left.text = "Technical Response"
    p_left.add_run(f"\n{project_name}").bold = True
    p_left.alignment = WD_ALIGN_PARAGRAPH.LEFT

    # Right cell for logo
    right_cell = table.cell(0, 1)
    right_cell.width = Inches(2.0)
    p_right = right_cell.paragraphs[0]
    try:
        p_right.add_run().add_picture(LOGO_PATH, width=Inches(1.5))
    except FileNotFoundError:
        print(f"INFO: Logo file not found at '{LOGO_PATH}'. The header will not contain a logo.")
        p_right.text = "[Company Logo]"
    p_right.alignment = WD_ALIGN_PARAGRAPH.RIGHT

def create_main_footer(section):
    """Creates the main footer for the document using a table for alignment."""
    footer = section.footer
    footer.is_linked_to_previous = False

    # Use a table to precisely control left and right alignment
    table = footer.add_table(rows=1, cols=2, width=Inches(6.5))
    
    # --- Left Cell: Company Info ---
    left_cell = table.cell(0, 0)
    p_left = left_cell.paragraphs[0]
    run_left = p_left.add_run("© Bahwan CyberTek")
    font_left = run_left.font
    font_left.size = Pt(8)
    font_left.color.rgb = DocxRGBColor(0x00, 0x00, 0xFF) # Blue
    p_left.alignment = WD_ALIGN_PARAGRAPH.LEFT

    # --- Right Cell: Page Number ---
    right_cell = table.cell(0, 1)
    p_right = right_cell.paragraphs[0]
    p_right.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # Add text "Page " before the page number field
    run_page_text = p_right.add_run("Page ")
    font_page_text = run_page_text.font
    font_page_text.color.rgb = DocxRGBColor(0x00, 0x00, 0xFF)
    font_page_text.size = Pt(8)

    # XML construction for the dynamic page number field
    field_begin = OxmlElement('w:fldChar')
    field_begin.set(qn('w:fldCharType'), 'begin')

    instr_text = OxmlElement('w:instrText')
    instr_text.set(qn('xml:space'), 'preserve')
    instr_text.text = r'PAGE \* MERGEFORMAT'

    field_end = OxmlElement('w:fldChar')
    field_end.set(qn('w:fldCharType'), 'end')

    # Add the XML elements to the paragraph's run
    p_right.runs[-1]._r.append(field_begin)
    p_right.runs[-1]._r.append(instr_text)
    p_right.runs[-1]._r.append(field_end)


def add_confidentiality_page(doc):
    """Adds the confidentiality page to the document."""
    doc.add_heading("Confidentiality", level=1)
    p_text = """The data in this document contains trade secrets and confidential or proprietary information of Bahwan CyberTek (BCT), the disclosure of which would provide a competitive advantage to others. As a result, this document shall not be disclosed, used or duplicated, in whole or in part, for any other purpose than to evaluate BCT. Disclosure of any data contained shall be with the express written permission of BCT. The data subject to this restriction are contained in the entire document.If a contract is awarded to BCT, as a result of, or in connection with the submission of this document, any right to duplicate, use, or disclose the data will be to the extent provided in the resulting agreement. This restriction does not limit the rights of the recipient to use the information contained in the data if it is rightfully obtained from another source without restriction. """
    p = doc.add_paragraph(p_text)
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    doc.add_page_break()

def add_table_of_contents(doc):
    """Adds a static, hierarchical table of contents without page numbers."""
    doc.add_heading("Table of Contents", level=1)

    # Define the structure of the TOC. 'level' controls the indentation.
    toc_structure = [
        {"level": 1, "text": "COVER LETTER"},
        {"level": 1, "text": "PROFILE OF FIRM"},
        {"level": 2, "text": "1.1. Bahwan CyberTek Overview and History"},
        {"level": 2, "text": "1.2. Financial Stability"},
        {"level": 1, "text": "WORK PLAN /PROJECT UNDERSTANDING"},
        {"level": 2, "text": "2.1. Our Understanding of VTA’s Objectives"},
        {"level": 2, "text": "2.2. Functionality Mapping"},
        {"level": 2, "text": "2.3. Architecture Design"},
        {"level": 2, "text": "2.4. Work Plan"},
        {"level": 1, "text": "APPENDIX A – BCT OFFERINGS, PARTNERSHIP & AWARDS,"},
        {"level": 2, "text": "3.1. Services Offering"},
        {"level": 2, "text": "3.2. Successful Strategic Partnerships and Joint Ventures"},
        {"level": 2, "text": "3.3. Awards and Recognition"},
        {"level": 2, "text": "3.4. BCT’s Experience – Public Sectors/Government"},
        {"level": 1, "text": "APPENDIX B: BCT CERTIFICATIONS"},
    ]
    
    for item in toc_structure:
        p = doc.add_paragraph()
        p.add_run(item["text"])
        
        # Apply indentation for hierarchical levels
        if item["level"] == 2:
            p.paragraph_format.left_indent = DocxInches(0.5)
            
    doc.add_page_break()


def add_cover_letter(doc):
    """Adds the hardcoded cover letter."""
    doc.add_heading("COVER LETTER", level=1)
    doc.add_paragraph("Dear Team,")
    
    p_sub = doc.add_paragraph()
    p_sub.add_run("Sub: Response to RFP for Utility Management System and Services (RFP S24094)").bold = True
    
    body = [
        "Our team at Bahwan CyberTek Inc (‘BCT’) has a long legacy of good work for entrepreneurs, governments, and large institutions. We have aided key initiatives for governments, world over, and are proud of having operated from the United States. The state of California is, to many of us, a marquee destination. In reviewing the requirements - we are certain we can make a difference to you and given more information – improve the strategic deployment of your assets, resources and more to meet your goals. BCT is excited to submit this proposal to Santa Clara Valley Transportation Authority (VTA) for the Utility Management System and Services.",
        "We understand that VTA is currently looking for a qualified firm to provide professional services for a Utility Management System and Services. After carefully reviewing the technical and operational requirements in detail, BCT proposes our own IP solution dtWorks. Our solution is specifically designed to meet the unique needs of utility management, offering robust features such as bill processing workflow, data visualization and seamless integration with existing system. We believe that dtWorks will fulfill the current requirements and also serve as a scalable platform to supports future growth. BCT brings the following capabilities and uniqueness for the current & future initiatives of VTA:",
    ]
    for para in body:
        doc.add_paragraph(para)
    
    bullet_points = [
        "Local presence in Santa Clara: BCT is in Santa Clara downtown and has its team of consultants located there. With a vast experience of working with local customers like the Housing Authority of City of Los Angeles we believe that BCT’s local presence will be of tremendous value for VTA.",
        "Transportation Domain Understanding: We are engaged in providing various implementation and support services for Rail, Road and Airlines industry. Some of our clients includes Dubai Metro, Roads and Transport Authority (RTA) of Dubai, COPA Airlines, Emirates Airlines Mumbai Metro Rail Corporation Limited (MMRCL-India), Autorent, Budget Rent-A-Car etc.",
        "SEI-CMMi Level 5 Assessment: BCT is proud to hold a SEI-CMMi Level 5 assessment, representing the highest level of process maturity. This distinction guarantees high-quality services, adherence to best practices, and predictable outcomes, reinforcing our commitment to excellence in strategic and operational IT needs.",
        "Experience working with Public Sectors: BCT has a strong presence in public sector/e-Government space and helped many cities worldwide to make a space for themselves in the list of the best e-governed cities of the world. BCT has played a pivotal role in large e-Government programs across multiple countries including the Housing Authority of City of Los Angeles, Houston PWE, City of Houston, Dubai Health Authority, Department of Health-Abu Dhabi etc.",
        "As evidence of our commitment to serve VTA, we are prepared to commit an effective team of highly skilled professionals and domain experts to support this endeavor."
    ]
    for point in bullet_points:
        doc.add_paragraph(point, style='List Bullet')


    doc.add_paragraph("\nSincerely,")
    doc.add_paragraph("Subramanian Nagarajan,")
    doc.add_paragraph("COO")
    doc.add_page_break()

def add_profile_of_firm(doc):
    """Adds the hardcoded Profile of Firm section."""
    doc.add_heading("PROFILE OF FIRM", level=1)
    doc.add_paragraph("In this section, we have provided our Company overview and history, organizational chart and financial stability in detail.")
    
    doc.add_heading("1.1. Bahwan CyberTek Overview and History", level=2)
    body1 = [
        "Bahwan CyberTek Group was founded in 1999, a strong multi-national company having a presence in over 50 countries across 4 continents. BCT is a recognized thought leader and innovative solutions partner for global Fortune 500 organizations and has delivered transformational solutions through IP-led products and cognitive solutions, growth accelerators and innovative outcome-based business models.",
        "BCT is a SEI CMMI Level 5 Company, offering the best IT processes in the industry. We are one of the few to have this distinction and have more than 4500+ employees working across the globe. This means that we are in a better position over others to provide the best and highest quality services to clients, as our processes comprehensively cover all aspects of implementation and support including risks mitigation and predictable project results.",
        "BCT has partnered with leading global technology organizations such as Oracle, Microsoft, IBM, TIBCO to deliver differentiated value to customers. BCT, today, has over 1000+ Enterprise Customers and 475+ Universities across the world including companies such as Department of Health Abu Dhabi, Dubai Health Authority, Cleveland clinic, UAB, Weaver, Viva Insurance, Housing Authority of the City of Los Angeles, UAB, Maurices, HomeServe, Biogen, Barnes & Noble, Nook, Tesla etc.",
        "BCT has operations in several countries including 5 offices in North America. This includes US Headquarters in Natick, MA; Tri-state Operations in Hackensack NJ, Southern Operations in Birmingham AL, West Coast Operations in Santa Clara CA and newly opened office in Texas. We are expanding in North America and are planning to open additional offices in South California in the near future."
    ]
    for para in body1:
        doc.add_paragraph(para)
        
    doc.add_heading("1.2. Financial Stability", level=2)
    body2 = [
        "BCT is an associate company of the Bahwan Group - a USD $12 billion conglomerate. Since our inception in 1999, BCT has maintained a consistent year-over-year growth record with revenues exceeding USD 300 million. BCT has been engaged with leading global corporations and established multiple, large scale, long-term contracts of sizes similar and / or larger. BCT has also made investments in high growth technology startups, that complement our ability to deliver extended value to our customers. This includes unicorns like Tekion, Cloudleaf in the IoT and SmarterD. With continued investments in its intellectual property, infrastructure capabilities and people, BCT maintains the leadership position in a number of areas and is poised for significant growth in the coming years.",
        "BCT hereby confirms that we have no pending litigation, and no claims or settlements have been paid by BCT or its insurers for any projects within the last five (5) year"
    ]
    for para in body2:
        doc.add_paragraph(para)

    doc.add_page_break()

def add_appendix_a(doc):
    """Adds the hardcoded Appendix A."""
    doc.add_heading("APPENDIX A – BCT OFFERINGS, PARTNERSHIP & AWARDS,", level=1)
    
    doc.add_heading("3.1. Services Offering", level=2)
    doc.add_paragraph("BCT’s services encompass the entire lifecycle of technologies and applications in the enterprise from Strategy to Support.")
    
    doc.add_heading("3.2. Successful Strategic Partnerships and Joint Ventures", level=2)
    doc.add_paragraph("BCT is Technology agnostic and works with the customer invested tools and technologies to maximize ROI. While we partner with most of the best of the breed vendors across the technology stack, we have a very customer focused approach and believe in value delivered to the customer than promoting our partnerships. Below are some of the tech partnerships active currently. We continue to evolve newer partnerships as needed for maximizing our value delivered.")
    
    p_d_offerings = doc.add_paragraph()
    p_d_offerings.add_run("BCT’s Digital Offerings").bold = True
    doc.add_paragraph("BCT is an IP-led enterprise that combines outcome-based business models, cognitive solutions and growth accelerators to deliver highly differentiated value to customers to help them achieve high levels of efficiency and performance. BCT offers digital solutions in key emerging areas that are seeing increasing demand from large organizations in different verticals.")
    
    digital_offerings = {
        "Digital Experience": [
            "DropThought is a comprehensive feedback management and intelligence engine that powers businesses to gather instant feedback from their customers, privately. dropthought uses customer sentiment analytics to provide powerful real-time insights. This allows customers and businesses to work together and improve experiences, interactions, events and overall trends",
            "CueTrack is a comprehensive customer service management product solution that helps organizations address and manage customer complaints and improve customer satisfaction. CueTrack provides a 360-degree view of customer cases to service desk officials. This helps improve transparency, reduce turnaround time and eliminate bottlenecks"
        ],
        "Digital Supply Chain Management": [
            "CueTrans An Integrated digital supply chain management suite of products with powerful HSE driven Journey Management and Track & Trace capabilities. CueTrans has a powerful and robust real-time fleet tracking and monitoring system with advanced tools for geofencing, driver behavior analysis, live monitoring and complete visibility over fleet’s movements and utilization.",
            "FuelTrans is an integrated logistics & transportation solution for Oil Marketing Companies, which offers a 360° view of the entire oil distribution logistics operations. FuelTrans enables Oil Marketing companies to manage and optimize their logistics assets and resources and brings full visibility and automation of the Transportation Process, from request to planning and fulfilment. This system also provides a Secondary Distribution Information Hub that gives a 360° view of oil distribution, logistics, and operations.",
            "Procure360 - An integrated AI-powered Supplier Management Solution that enables digital transformation across your entire procurement and vendor value chain, from source to pay"
        ],
        "Predictive Analytics": [
            "Retina360 - Retina360 is an integrated analytics platform with powerful capabilities to collect, validate and process data from multiple data sources and provide multi-parameter predictive modelling, using advanced statistical methods, AI techniques, evolutionary algorithms and constraint optimization tools. Solutions developed on the Retina360 platform work in asset-intensive organizations, helping monitor, maintain and optimize assets for better utilization and performance.",
            "Rt360 – rt360 is a suite of Risk Management solutions for the BFSI space, commonly known as Regulatory Technology. The suite includes solutions for the management of credit risk, liquidity risk, operational risk and model risk in the banking context.",
            "Geodatafy – an Exploration & Production digitalization platform that provides complete transparency from field to operations through integration of technical data across Oil and Gas business units, such as Drilling, Subsurface, Operations, Finance and HSSE."
        ]
    }
    for category, items in digital_offerings.items():
        doc.add_paragraph().add_run(category).bold = True
        for item in items:
            doc.add_paragraph(item, style='List Bullet')

    doc.add_heading("3.3. Awards and Recognition", level=2)
    doc.add_paragraph("Over the years, BCT has been recognized by global industry bodies, partners and customers with prestigious awards.")
    p_stanford = doc.add_paragraph()
    p_stanford.add_run("Stanford Case Study").bold = True
    doc.add_paragraph("As testimony to its journey over the years, Stanford Graduate School of Business has written two case studies on the evolution of the Bahwan CyberTek Group. The case studies written in 2017 & 2018 are unique in that Bahwan CyberTek is the only company from the region to be featured by Stanford University. Going further BCT’s case study has been included in the curriculum for Management Studies in Stanford in Palo Alto USA and also listed in HBS.")
    
    p_awards = doc.add_paragraph()
    p_awards.add_run("Awards").bold = True
    awards_list = [
        "Skoch Award for Urban Development, 2023 – BCT’s client, the Directorate of Municipal Administration, won this award.",
        "BCT TIBCO was one of the winners at The Economic Times Choice of Tech Leaders 2023 for \"Product & Services Strategy\"",
        "Excellence in Asset Management & Maintenance Award, 2023 - BCT’s Enterprise Asset Management (EAM) team was recognized with the Excellence in Asset Management and Maintenance Award at the 4th Rail Analysis Innovation & Excellence Summit 2023",
        "Outlook Business (Business Icon & Brand of the Year) Award 2022 – Outlook Business recognized Mr. S. Durgaprasad (DP) for his exemplary leadership and BCT for its exceptional commitment to product development.",
        "Strategic Technical Partner Award at GITEX Global 2022 - BCT was bestowed with the coveted Strategic Technical Partner Award, presented by Ajman Municipality and Planning Department (AMPD).",
        "FICCI Runner-Up Award for Best R&D and Product Development - retina360, a BCT product, won the prestigious FICCI Runner-up award for the \"Best R&D and Product Development” at the 12th edition of the Tanenergy Summit & Awards 2022, organized by The Federation of Indian Chambers of Commerce & Industry (FICCI) Tamil Nadu State Council",
        "Sonatype Partner Award 2022 - BCT won the Sonatype EMEA Partner Awards 2022",
        "BCT is a Microsoft Gold Partner - BCT achieved the Gold Partner status with Microsoft for the 9th Consecutive Year",
        "Multi-Sectoral Digital Transformation Award at OER Live 2022, Oman - BCT received the award for ‘Achievements in Multi-Sectoral Digital Transformation’",
        "Sir Viswesvaraya Award for 'Best Corporate Entity' from AIMO at TechKnow'22 - BCT was presented with the Sir M. Visvesvaraya Best Corporate Entity Award by the Chief Minister of Tamil Nadu, Mr. MK Stalin at TechKnow 2022",
        "FuelTrans won the Confederation of Indian Industries (CII) SCALE 2022 Award for the best “Outstanding Performance” in the supply chain and logistics domain.",
        "Asia Pacific Enterprise Award for Corporate Excellence, 2021",
        "Best Service Provider of the Year for Data Analytics (Wind), 2021",
        "‘Outstanding Employee Engagement Strategy Award’ and ‘Business HR Award’, 2021",
        "Best Talent Development Award & the Learning & Development Collaboration Award, 2021",
        "rt360 recognized as ‘Rising Star of the Year’ at the RiskTech100® 2021.",
        "FuelTrans wins the prestigious CII SCALE Award for Best Technology Solution Provider, 2021"
    ]
    for award in awards_list:
        doc.add_paragraph(award, style='List Bullet')

    doc.add_heading("3.4. BCT’s Experience – Public Sectors/Government", level=2)
    doc.add_paragraph("BCT has a strong presence in public sector/e-Government space and helped many cities worldwide to make a space for themselves in the list of the best e-governed cities of the world. BCT has played a pivotal role in large e-Government programs across multiple countries offering e-Payment, Integrated Payment Hub, Payment Aggregation, Integrated Revenue Management, Business Licensing, Citizen Services, Managed services of citizen contact center. With rich experience in implementation, upgrade and support our specialists conceptualize, develop, deliver and operate the best-in-class governance solutions which promote efficiency, transparency and enable “smart governance”.")
    
    experience_data = [
        ("BCT is engaged with HACLA to provide application managed services which includes core application development, application enhancement, cross functional services, service desk, infrastructure services, network and end user computing."),
        ("BCT assisted Houston PWE and has worked with them since 2011 in architecture, installation and configuration, design, development, deployment and support of their core enterprise systems such as utility billing, GIS, and other applications by establishing a ‘standards-based Middleware’ platform and integrated the systems in a seamless manner"),
        ("BCT engaged with NY Department of Health in delivering enterprise integration services for their core enterprise applications"),
        ("BCT assisted in the migration of database for Cook County Hospital’s Cerner Hospital Management System."),
        ("BCT has been engaged with University of Alabama Birmingham School of Medicine for many years and in a number of critical IT initiatives including application and infrastructure managed services, development, testing and quality assurance, production support, etc."),
        ("BCT assisted the City of Chandler, AZ in designing an enterprise integration architecture for seamlessly integrating all the City’s Information Systems using a standards-based middleware platform.")
    ]
    for item in experience_data:
        doc.add_paragraph(item, style='List Bullet')

    doc.add_page_break()

def add_appendix_b(doc):
    """Adds the hardcoded Appendix B."""
    doc.add_heading("APPENDIX B: BCT CERTIFICATIONS", level=1)
    doc.add_paragraph("In this section, we have provided the list of certifications.")
    certifications = [
        "ISO/IEC 20000-1:2018",
        "ISO 9001:2015",
        "ISO 27001:2013",
        "SSAE 18 SOC TYPE II",
        "CMMI Level 5"
    ]
    for cert in certifications:
        doc.add_paragraph(cert, style='List Bullet')
    doc.add_page_break()


def generate_technical_proposal_docx(understanding_data, rfp_data):
    """Generates the 'Technical Proposal' DOCX file based on LLM analysis and hardcoded sections."""
    doc = docx.Document()
    
    # Set default font for the document
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Arial'
    font.size = Pt(11)
    
    # --- Section and Header/Footer Setup ---
    project_name_val = rfp_data.get("Prospect & RFP Background", {}).get("Project Name", "Project Name")
    section = doc.sections[0]
    section.different_first_page_header_footer = True
    create_main_header(section, project_name_val)
    create_main_footer(section)

    # --- 1. Title Page ---
    client_name_val = rfp_data.get("Prospect & RFP Background", {}).get("Client Name", "N/A")
    rfp_number_val = rfp_data.get("Prospect & RFP Background", {}).get("RFP Number", "N/A")
    submission_date = date.today().strftime("%d %b, %Y")

    def add_centered_paragraph(text, bold=False, font_size=12):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.bold = bold
        run.font.size = DocxPt(font_size)
        return p

    add_centered_paragraph("Bahwan CyberTek – Technical Response to the", font_size=14)
    add_centered_paragraph("Request for Proposals for", font_size=14)
    doc.add_paragraph()
    add_centered_paragraph(project_name_val, bold=True, font_size=18)
    add_centered_paragraph(rfp_number_val, font_size=14)
    doc.add_paragraph()
    doc.add_paragraph()
    add_centered_paragraph("Submitted to", font_size=14)
    add_centered_paragraph(client_name_val, bold=True, font_size=16)
    doc.add_paragraph()
    doc.add_paragraph()
    add_centered_paragraph("Submitted On", font_size=14)
    add_centered_paragraph(submission_date, bold=True, font_size=14)
    doc.add_page_break()
    
    # --- 2. Confidentiality Page ---
    add_confidentiality_page(doc)

    # --- 3. Table of Contents (MODIFIED) ---
    add_table_of_contents(doc)

    # --- 4. Main Content Sections ---
    add_cover_letter(doc)
    add_profile_of_firm(doc)
    
    # --- WORK PLAN /PROJECT UNDERSTANDING (LLM Generated) ---
    doc.add_heading("WORK PLAN /PROJECT UNDERSTANDING", level=1)

    # 2.1 Understanding Objectives
    if "Understanding Objectives" in understanding_data:
        doc.add_heading("2.1. Our Understanding of VTA’s Objectives", level=2)
        obj_data = understanding_data["Understanding Objectives"]
        # This will render all sub-sections of the LLM's understanding objectives
        if "Strategic Objectives" in obj_data and obj_data["Strategic Objectives"]:
            for item in obj_data["Strategic Objectives"]:
                if str(item).strip():
                    doc.add_paragraph(str(item).strip(), style='List Bullet')
    
    # 2.2 Functionality Mapping
    if "Functionality Mapping" in understanding_data:
        doc.add_heading("2.2. Functionality Mapping", level=2)
        content_map = understanding_data["Functionality Mapping"]
        if isinstance(content_map, dict) and "functionalities" in content_map:
            for func_item in content_map["functionalities"]:
                if isinstance(func_item, dict) and func_item.get("name"):
                    doc.add_paragraph().add_run(func_item["name"]).bold = True
                    for key, value in func_item.items():
                        if key != "name" and value:
                            p = doc.add_paragraph(style='List Bullet')
                            p.add_run(f"{key.replace('_', ' ').title()}:").bold = True
                            p.add_run(f" {value}")
    
    # 2.3 Architecture Design
    if "Architecture Design" in understanding_data:
        doc.add_heading("2.3. Architecture Design", level=2)
        arch_data = understanding_data["Architecture Design"]
        if isinstance(arch_data, dict) and "components" in arch_data:
                      for comp_item in arch_data["components"]:
                          if isinstance(comp_item, dict) and comp_item.get("name"):
                              doc.add_paragraph().add_run(comp_item["name"]).bold = True
                              for key, value in comp_item.items():
                                  if key != "name" and value:
                                      p = doc.add_paragraph(style='List Bullet')
                                      p.add_run(f"{key.replace('_', ' ').title()}:").bold = True
                                      p.add_run(f" {value}")

    # 2.4 Work Plan
    if "Work Plan (Detail)" in understanding_data:
        doc.add_heading("2.4. Work Plan", level=2)
        work_plan_data = understanding_data["Work Plan (Detail)"]
        if isinstance(work_plan_data, dict) and "phases" in work_plan_data:
                      for phase_item in work_plan_data["phases"]:
                          if isinstance(phase_item, dict) and phase_item.get("name"):
                              doc.add_paragraph().add_run(phase_item["name"]).bold = True
                              for key, value in phase_item.items():
                                  if key != "name" and value:
                                      p = doc.add_paragraph(style='List Bullet')
                                      p.add_run(f"{key.replace('_', ' ').title()}:").bold = True
                                      p.add_run(f" {value}")
    doc.add_page_break()

    # --- Remaining Hardcoded Sections ---
    # The "ADMINISTRATIVE SUBMITTALS" section is now removed.
    
    add_appendix_a(doc)
    add_appendix_b(doc)


    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# --- Streamlit Application ---

st.set_page_config(layout="wide", page_title="RFP Document Chatbot")

# --- Initialize Session State ---
if "messages" not in st.session_state:
    st.session_state.messages = []
if "processed_doc_content" not in st.session_state:
    st.session_state.processed_doc_content = None
if "groq_model_name" not in st.session_state:
    st.session_state.groq_model_name = "deepseek-r1-distill-llama-70b"
if "chat_model_name" not in st.session_state:
    st.session_state.chat_model_name = "llama-3.3-70b-versatile"
if "current_uploaded_file_name" not in st.session_state:
    st.session_state.current_uploaded_file_name = None
if "processing_done" not in st.session_state:
    st.session_state.processing_done = False
if "doc_analysis_data" not in st.session_state:
    st.session_state.doc_analysis_data = {}
if "understanding_doc_analysis_data" not in st.session_state:
    st.session_state.understanding_doc_analysis_data = {}
if "title_bg_image_bytes" not in st.session_state:
    st.session_state.title_bg_image_bytes = None
if "title_bg_image_name" not in st.session_state:
    st.session_state.title_bg_image_name = None
if "second_page_bg_image_bytes" not in st.session_state:
    st.session_state.second_page_bg_image_bytes = None
if "second_page_bg_image_name" not in st.session_state:
    st.session_state.second_page_bg_image_name = None
if "content_bg_image_bytes" not in st.session_state:
    st.session_state.content_bg_image_bytes = None
if "content_bg_image_name" not in st.session_state:
    st.session_state.content_bg_image_name = None

# --- FIX: Check if processing is done and reload vector store if needed ---
# This ensures that the vector store is available across app reruns.
if st.session_state.processing_done:
    try:
        initialize_components()
    except Exception as e:
        st.error(f"Could not reload the vector store. Please re-upload your document. Error: {e}")
        st.session_state.processing_done = False


# --- Custom CSS for better chat display ---
st.markdown("""
<style>
.st-chat-message-contents {
    background-color: #f0f2f6;
    border-radius: 10px;
    padding: 10px;
}
.st-chat-message-contents.ai {
    background-color: #e6f7ff;
}
</style>
""", unsafe_allow_html=True)


st.title("BCT BID Management")
st.write("Upload your RFP document (PDF, DOCX, PPTX) to chat with it!")

# --- Sidebar for controls and info ---
with st.sidebar:
    st.header("Settings")
    st.info("Your document will be processed and indexed in a vector database for efficient querying. The chatbot uses Groq for fast inference.")

    st.success(f"LLM for analysis: **{st.session_state.groq_model_name}**")
    st.success(f"LLM for chat: **{st.session_state.chat_model_name}**")


    st.markdown("---")
    st.header("PPTX Design Settings")
    st.write("Upload images to use as backgrounds for the generated PPTX slides.")

    # Uploader for Title Slide Background
    uploaded_title_bg_image = st.file_uploader(
        "**1. Choose Title Slide Background (PNG, JPG, JPEG)**",
        type=["png", "jpg", "jpeg"],
        key="title_bg_image_uploader"
    )
    if uploaded_title_bg_image is not None:
        st.session_state.title_bg_image_bytes = uploaded_title_bg_image.getvalue()
        st.session_state.title_bg_image_name = uploaded_title_bg_image.name
        st.success(f"Title slide background loaded: {uploaded_title_bg_image.name}")
    else:
        st.session_state.title_bg_image_bytes = None
        st.session_state.title_bg_image_name = None
        st.info("No title slide background uploaded. First slide will have a default white background.")

    # Uploader for Second Page Background
    uploaded_second_page_bg_image = st.file_uploader(
        "**2. Choose Second Page Background (PNG, JPG, JPEG)**",
        type=["png", "jpg", "jpeg"],
        key="second_page_bg_image_uploader"
    )
    if uploaded_second_page_bg_image is not None:
        st.session_state.second_page_bg_image_bytes = uploaded_second_page_bg_image.getvalue()
        st.session_state.second_page_bg_image_name = uploaded_second_page_bg_image.name
        st.success(f"Second page background loaded: {uploaded_second_page_bg_image.name}")
    else:
        st.session_state.second_page_bg_image_bytes = None
        st.session_state.second_page_bg_image_name = None
        st.info("No specific background for the second page. It will use the default or the 'Subsequent Pages' background if uploaded.")

    # Uploader for Content Slides Background (from page 3 onwards)
    uploaded_content_bg_image = st.file_uploader(
        "**3. Choose Subsequent Pages Background (PNG, JPG, JPEG)**",
        type=["png", "jpg", "jpeg"],
        key="content_bg_image_uploader"
    )
    if uploaded_content_bg_image is not None:
        st.session_state.content_bg_image_bytes = uploaded_content_bg_image.getvalue()
        st.session_state.content_bg_image_name = uploaded_content_bg_image.name
        st.success(f"Subsequent pages background loaded: {uploaded_content_bg_image.name}")
    else:
        st.session_state.content_bg_image_bytes = None
        st.session_state.content_bg_image_name = None
        st.info("No background for subsequent pages. They will have a default white background.")


# --- Main area for document upload and processing ---
st.subheader("1. Upload RFP Document")
uploaded_file = st.file_uploader(
    "Choose a file",
    type=["pdf", "docx", "pptx"],
    help="Upload your RFP document. Supported formats: PDF, DOCX, PPTX.",
    key="file_uploader"  
)

# Logic to handle file upload and processing
if uploaded_file is not None:
    if uploaded_file.name != st.session_state.current_uploaded_file_name or not st.session_state.processing_done:
        st.session_state.messages = []
        
        with tempfile.TemporaryDirectory() as tmpdir:
            temp_file_path = Path(tmpdir) / uploaded_file.name
            with open(temp_file_path, "wb") as f:
                f.write(uploaded_file.getvalue())

            file_type = uploaded_file.type.split('/')[-1]

            st.info(f"Processing your **{file_type.upper()}** file: **{uploaded_file.name}**")
            
            processing_success = process_document_and_update_state(temp_file_path, file_type, uploaded_file.name)

            if processing_success:
                st.success("Document processing complete! You can now chat with your document.")
                
                # Run LLM analysis for document generation after successful text extraction
                with st.spinner("Performing advanced document analysis for PPTX/DOCX generation..."):
                    st.session_state.doc_analysis_data = analyze_with_llm_for_docs(
                        st.session_state.processed_doc_content, # type: ignore
                        st.session_state.groq_model_name # Use the hardcoded model name
                    )
                if st.session_state.doc_analysis_data:
                    st.success("Structured analysis for DOCX/PPTX generation complete!")
                else:
                    st.warning("Could not complete structured analysis for DOCX/PPTX generation. Generated documents might be empty or incomplete.")

                # NEW: Run LLM analysis for the "Understanding Doc"
                with st.spinner("Generating 'Technical Proposal' content with advanced LLM analysis..."):
                    st.session_state.understanding_doc_analysis_data = analyze_for_understanding_doc(
                        st.session_state.processed_doc_content, # type: ignore
                        st.session_state.groq_model_name
                    )
                if st.session_state.understanding_doc_analysis_data:
                    st.success("'Technical Proposal' content generated!")
                else:
                    st.warning("Could not generate 'Technical Proposal' content. The document might be empty or incomplete.")


            else:
                st.warning("Document processing might not have completed successfully. Please check the console/messages above for details.")
                st.session_state.processed_doc_content = None
                st.session_state.current_uploaded_file_name = None
                st.session_state.processing_done = False
                st.session_state.understanding_doc_analysis_data = {} # Clear this too on failure
    elif uploaded_file.name == st.session_state.current_uploaded_file_name and st.session_state.processing_done:
        st.info(f"Currently loaded document: **{st.session_state.current_uploaded_file_name}**")
else:
    st.session_state.processed_doc_content = None
    st.session_state.current_uploaded_file_name = None
    st.session_state.processing_done = False
    st.session_state.doc_analysis_data = {} # Ensure analysis data is cleared
    st.session_state.understanding_doc_analysis_data = {} # Clear analysis for Understanding Doc


# --- 2. Download Generated Documents ---
st.subheader("2. Download Generated Documents")
if st.session_state.processed_doc_content and st.session_state.current_uploaded_file_name:
    st.write("You can download the processed content in different formats:")

    col1, col2, col3 = st.columns(3)

    with col1:
        if st.session_state.doc_analysis_data:
            docx_file_bytes = generate_docx_option1(st.session_state.doc_analysis_data)
            st.download_button(
                label="Download DOCX (RFP Analysis)",
                data=docx_file_bytes,
                file_name=f"{st.session_state.current_uploaded_file_name}_RFP_Analysis.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                key="download_detailed_docx"
            )
        else:
            st.button("Download DOCX (RFP Analysis)", disabled=True, help="Complete LLM analysis first.")

    with col2:
        if st.session_state.doc_analysis_data and \
           (st.session_state.title_bg_image_bytes or st.session_state.second_page_bg_image_bytes or st.session_state.content_bg_image_bytes):
            
            title_bg_info = {
                "bytes": st.session_state.title_bg_image_bytes,
                "name": st.session_state.title_bg_image_name or "default_title_bg.png"
            }
            second_page_bg_info = {
                "bytes": st.session_state.second_page_bg_image_bytes,
                "name": st.session_state.second_page_bg_image_name or "default_second_page_bg.png"
            }
            content_bg_info = {
                "bytes": st.session_state.content_bg_image_bytes,
                "name": st.session_state.content_bg_image_name or "default_content_bg.png"
            }

            generated_pptx_bytes = generate_new_pptx_option1(
                analysis_data=st.session_state.doc_analysis_data,
                title_bg_info=title_bg_info,
                second_page_bg_info=second_page_bg_info,
                content_bg_info=content_bg_info
            )
            st.download_button(
                label="Download PPTX (RFP Analysis)",
                data=generated_pptx_bytes,
                file_name=f"{st.session_state.current_uploaded_file_name}_RFP_Analysis.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="download_detailed_pptx_with_bg"
            )
        else:
            st.info("Upload a document and at least one background image to enable this PPTX download.")
            
    with col3:
        if st.session_state.understanding_doc_analysis_data and st.session_state.doc_analysis_data:
            technical_proposal_bytes = generate_technical_proposal_docx(
                st.session_state.understanding_doc_analysis_data,
                st.session_state.doc_analysis_data
            )
            st.download_button(
                label="Download Technical Proposal",
                data=technical_proposal_bytes,
                file_name=f"{st.session_state.current_uploaded_file_name}_Technical_Proposal.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                key="download_technical_proposal"
            )
        else:
            st.button("Download Technical Proposal", disabled=True, help="Upload a document and wait for all analyses to complete.")

else:
    st.info("Upload a document first to enable document generation options.")


# --- 3. Chat with your Document ---
st.subheader("3. Chat with your Document")
for message in st.session_state.messages:
    with st.chat_message(message["role"]):
        st.markdown(message["content"])

if query := st.chat_input("Ask a question about the document..."):
    st.chat_message("user").markdown(query)
    st.session_state.messages.append({"role": "user", "content": query})
    with st.spinner("Thinking..."):
        try:
            answer, sources = generate_answer(query)
            st.session_state.messages.append({"role": "assistant", "content": answer})
            with st.chat_message("assistant"):
                st.markdown(answer)
                if sources:
                    unique_sources = list(set(sources))
                    st.markdown(f"**Sources:** {', '.join(unique_sources)}")
        except RuntimeError as e:
            st.error(str(e))
        except Exception as e:
            st.error(f"An error occurred during chat: {e}")
